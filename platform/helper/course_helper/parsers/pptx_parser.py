"""Deterministic, notes-first extraction from registered PPTX sources."""

from __future__ import annotations

import hashlib
import json
import re
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Any
from uuid import uuid5

from PIL import Image
from pptx import Presentation, __version__ as PPTX_VERSION
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from course_helper.domain.common import ActorRef, SourceLocator
from course_helper.domain.evidence import EvidenceCheck, EvidenceObject
from course_helper.domain.sources import (
    ChunkLocator,
    ExtractedChunk,
    ExtractionResult,
    SourceAssetVersion,
    VisualAssetVersion,
)
from course_helper.source_roots import (
    COURSE_STUDIO_ID_NAMESPACE,
    SourceRootRegistry,
    candidate_logical_id,
    candidate_version_id,
    chunk_logical_id,
    chunk_version_id,
    source_logical_id,
    source_version_id,
    stream_sha256,
)


PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PARSER_NAME = "python-pptx"
PARSER_PRODUCER = "course-helper/pptx-parser"
_ACTOR = ActorRef(actor_type="service", actor_id=PARSER_PRODUCER)
_RELATIONSHIP_NUMBER = re.compile(r"^(.*?)(\d+)$")
_BLOB_BLOCK_SIZE = 1024 * 1024


def normalize_text(value: str) -> str:
    """Normalize line endings and trailing whitespace without flattening structure."""

    lines = value.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    return "\n".join(line.rstrip() for line in lines).strip()


class PptxParser:
    """Parse allowlisted PPTX files into immutable extraction contracts."""

    def __init__(self, source_roots: SourceRootRegistry) -> None:
        self._source_roots = source_roots

    def parse(
        self,
        locator: SourceLocator,
        slide_range: range | None = None,
    ) -> ExtractionResult:
        path = self._source_roots.resolve(locator)
        source_digest = stream_sha256(path)
        source_logical = source_logical_id(locator)
        source_version = source_version_id(source_logical, source_digest)
        presentation = Presentation(path)
        selected_slides = tuple(
            (slide_number, slide)
            for slide_number, slide in enumerate(presentation.slides, start=1)
            if slide_range is None or slide_number in slide_range
        )
        selected_numbers = tuple(slide_number for slide_number, _ in selected_slides)
        parser_config_digest = _parser_config_digest(selected_numbers)
        source_time = _source_time(path)

        chunks: list[ExtractedChunk] = []
        visuals: list[VisualAssetVersion] = []
        checks: list[EvidenceCheck] = []
        failed_relationships = 0

        for ordinal, (slide_number, slide) in enumerate(selected_slides):
            slide_text = _slide_text(slide)
            notes_text = _notes_text(slide)
            normalized = normalize_text(
                "\n\n".join(part for part in (notes_text, slide_text) if part)
            )
            chunk_locator: ChunkLocator | None = None
            chunk_digest: str | None = None
            chunk_id: str | None = None
            if normalized:
                chunk_locator = ChunkLocator(kind="pptx-slide", slide_number=slide_number)
                chunk_digest = hashlib.sha256(normalized.encode("utf-8")).hexdigest()
                chunk_logical = chunk_logical_id(source_logical, chunk_locator)
                chunk_id = chunk_version_id(chunk_logical, source_version, chunk_digest)
            else:
                checks.append(
                    EvidenceCheck(
                        code="pptx-slide-text",
                        status="failed",
                        message=f"Slide {slide_number} has no extractable notes or slide text",
                        details={"slide_number": slide_number},
                    )
                )
            slide_visuals: list[VisualAssetVersion] = []
            warnings: list[str] = []

            image_relationships = sorted(
                (
                    relationship
                    for relationship in slide.part.rels.values()
                    if relationship.reltype == RT.IMAGE
                ),
                key=lambda relationship: _relationship_sort_key(relationship.rId),
            )
            for relationship in image_relationships:
                try:
                    visual = _visual_from_relationship(
                        relationship=relationship,
                        slide_number=slide_number,
                        source_logical_id_value=source_logical,
                        parent_version_ids=(
                            (source_version, chunk_id)
                            if chunk_id is not None
                            else (source_version,)
                        ),
                        created_at=source_time,
                    )
                except (
                    AssertionError,
                    AttributeError,
                    OSError,
                    TypeError,
                    ValueError,
                ) as error:
                    failed_relationships += 1
                    warnings.append(
                        f"Image relationship {relationship.rId} could not be extracted"
                    )
                    checks.append(
                        EvidenceCheck(
                            code="pptx-image-relationship",
                            status="failed",
                            message=(
                                f"Slide {slide_number} image relationship "
                                f"{relationship.rId} could not be extracted"
                            ),
                            details={
                                "slide_number": slide_number,
                                "relationship_id": relationship.rId,
                                "exception_type": type(error).__name__,
                            },
                        )
                    )
                    continue

                slide_visuals.append(visual)
                visuals.append(visual)
                checks.append(
                    EvidenceCheck(
                        code="pptx-image-relationship",
                        status="passed",
                        message=(
                            f"Slide {slide_number} image relationship "
                            f"{relationship.rId} was hashed and inspected"
                        ),
                        details={
                            "slide_number": slide_number,
                            "relationship_id": relationship.rId,
                            "content_digest": visual.content_digest,
                            "width": visual.width,
                            "height": visual.height,
                        },
                    )
                )

            if chunk_id is None or chunk_locator is None or chunk_digest is None:
                continue
            heading = next((line.strip() for line in slide_text.splitlines() if line.strip()), None)
            chunks.append(
                ExtractedChunk(
                    chunk_id=chunk_id,
                    source_version_id=source_version,
                    ordinal=ordinal,
                    modality="slide",
                    language="und",
                    normalized_text=normalized,
                    content_digest=chunk_digest,
                    locator=chunk_locator,
                    breadcrumb=(heading,) if heading is not None else (),
                    heading=heading,
                    notes_text=notes_text,
                    slide_text=slide_text,
                    media_version_ids=tuple(visual.version_id for visual in slide_visuals),
                    warnings=tuple(warnings),
                )
            )

        failed_checks = sum(check.status == "failed" for check in checks)
        extraction_status = "partial" if failed_checks else "parsed"
        evidence_status = "degraded" if failed_checks else "verified"
        source = SourceAssetVersion(
            logical_id=source_logical,
            version_id=source_version,
            revision=1,
            content_digest=source_digest,
            created_at=source_time,
            created_by=_ACTOR,
            locator=locator,
            display_name=Path(locator.relative_path).name,
            source_kind="pptx",
            media_type=PPTX_MEDIA_TYPE,
            byte_size=path.stat().st_size,
            modified_at=source_time,
            content_summary=f"{len(chunks)} slide chunks and {len(visuals)} visual assets",
            extraction_status=extraction_status,
            parser_name=PARSER_NAME,
            parser_version=PPTX_VERSION,
            parser_config_digest=parser_config_digest,
        )
        checks.insert(
            0,
            EvidenceCheck(
                code="pptx-extraction",
                status="warning" if failed_checks else "passed",
                message=(
                    f"Extracted {len(chunks)} selected slides and {len(visuals)} visual assets"
                ),
                details={
                    "selected_slide_count": len(selected_numbers),
                    "chunk_count": len(chunks),
                    "visual_count": len(visuals),
                    "failed_relationship_count": failed_relationships,
                },
            ),
        )
        evidence = EvidenceObject(
            evidence_id=_evidence_id(source_version, parser_config_digest),
            kind="extraction",
            subject_version_id=source_version,
            status=evidence_status,
            input_summary={
                "source_locator": locator.model_dump(mode="json"),
                "selected_slide_numbers": list(selected_numbers),
            },
            output_summary={
                "chunk_count": len(chunks),
                "visual_count": len(visuals),
                "failed_check_count": failed_checks,
            },
            producer=PARSER_PRODUCER,
            producer_version="1",
            started_at=source_time,
            finished_at=source_time,
            duration_ms=0,
            checks=tuple(checks),
        )
        return ExtractionResult(
            source=source,
            chunks=tuple(chunks),
            visuals=tuple(visuals),
            evidence=evidence,
        )


def _slide_text(slide: Any) -> str:
    return "\n".join(
        shape.text.strip()
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def _notes_text(slide: Any) -> str:
    if not slide.has_notes_slide:
        return ""
    notes_frame = slide.notes_slide.notes_text_frame
    return notes_frame.text.strip() if notes_frame is not None else ""


def _visual_from_relationship(
    *,
    relationship: Any,
    slide_number: int,
    source_logical_id_value: str,
    parent_version_ids: tuple[str, ...],
    created_at: datetime,
) -> VisualAssetVersion:
    if relationship.is_external:
        raise ValueError("external image relationship is not an embedded media part")
    target_part = relationship.target_part
    blob = target_part.blob
    content_digest = _streaming_blob_digest(blob)
    with Image.open(BytesIO(blob)) as image:
        width, height = image.size
        image.verify()

    semantic_locator = (
        f"{source_logical_id_value}:slide-{slide_number}:{relationship.rId}"
    )
    logical_id = candidate_logical_id("visual", semantic_locator)
    version_id = candidate_version_id(
        logical_id,
        parent_version_ids,
        content_digest,
    )
    return VisualAssetVersion(
        logical_id=logical_id,
        version_id=version_id,
        revision=1,
        content_digest=content_digest,
        created_at=created_at,
        created_by=_ACTOR,
        media_type=target_part.content_type,
        width=width,
        height=height,
        source_locator=ChunkLocator(
            kind="pptx-slide",
            slide_number=slide_number,
            relationship_id=relationship.rId,
        ),
        license_status="source-provided",
        authenticity="source-provided",
        derived_from_version_ids=parent_version_ids,
    )


def _streaming_blob_digest(blob: bytes) -> str:
    digest = hashlib.sha256()
    view = memoryview(blob)
    for offset in range(0, len(view), _BLOB_BLOCK_SIZE):
        digest.update(view[offset : offset + _BLOB_BLOCK_SIZE])
    return digest.hexdigest()


def _parser_config_digest(selected_slide_numbers: tuple[int, ...]) -> str:
    payload = json.dumps(
        {
            "parser": PARSER_NAME,
            "parser_version": PPTX_VERSION,
            "selected_slide_numbers": selected_slide_numbers,
        },
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _source_time(path: Path) -> datetime:
    return datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc)


def _evidence_id(source_version_id_value: str, parser_config_digest: str) -> str:
    return str(
        uuid5(
            COURSE_STUDIO_ID_NAMESPACE,
            f"evidence\0pptx\0{source_version_id_value}\0{parser_config_digest}",
        )
    )


def _relationship_sort_key(relationship_id: str) -> tuple[str, int, str]:
    match = _RELATIONSHIP_NUMBER.match(relationship_id)
    if match is None:
        return relationship_id, -1, relationship_id
    return match.group(1), int(match.group(2)), relationship_id
