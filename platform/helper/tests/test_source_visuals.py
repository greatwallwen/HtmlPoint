from __future__ import annotations

from datetime import datetime, timezone
import hashlib
from io import BytesIO
from pathlib import Path
import sqlite3
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from course_helper.catalog import KnowledgeCatalog
from course_helper.domain.common import SourceLocator
from course_helper.domain.sources import VisualAssetVersion
from course_helper.parsers.pptx_parser import PptxParser
from course_helper.source_roots import SourceRootRegistry, stream_sha256


NOW = datetime(2026, 7, 18, 12, 0, tzinfo=timezone.utc)


def png_bytes(color: tuple[int, int, int]) -> bytes:
    stream = BytesIO()
    Image.new("RGB", (4, 3), color=color).save(stream, format="PNG")
    return stream.getvalue()


def build_pptx(path: Path, images: tuple[bytes, ...]) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(5), Inches(0.5))
    box.text = "Traceable source visuals"
    for index, payload in enumerate(images):
        slide.shapes.add_picture(
            BytesIO(payload),
            Inches(0.5 + index * 1.25),
            Inches(1),
            Inches(1),
            Inches(0.75),
        )
    slide.notes_slide.notes_text_frame.text = "Explain only what the source shows."
    presentation.save(path)
    return path


def persist_extraction(catalog: KnowledgeCatalog, result) -> None:
    catalog.insert_source(result.source)
    for chunk in result.chunks:
        catalog.insert_chunk(chunk)
    for visual in result.visuals:
        catalog.insert_visual(visual)
    catalog.insert_evidence(result.evidence)


def test_materializes_exact_relationship_without_modifying_source(tmp_path: Path) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "source.pptx", (png_bytes((20, 40, 60)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    before_digest = stream_sha256(pptx)
    before_stat = pptx.stat()
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        outcomes = materialize_source_visuals(
            catalog,
            registry,
            ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[0].version_id,),
            clock=lambda: NOW,
        )
        assert len(outcomes) == 1
        outcome = outcomes[0]
        assert outcome.status == "materialized"
        assert outcome.artifact_id
        assert outcome.evidence_id
        assert outcome.error_code is None
        binding = catalog.get_source_visual_materialization(result.visuals[0].version_id)
        assert binding is not None
        assert binding.payload.artifact_id == outcome.artifact_id
        artifact = catalog.get_artifact(outcome.artifact_id)
        assert artifact is not None
        assert artifact.payload.content_digest == result.visuals[0].content_digest
        assert (artifact.payload.width, artifact.payload.height) == (4, 3)
        assert "path" not in binding.payload.model_dump_json()
        evidence_payload = catalog.connection.execute(
            "SELECT payload_json FROM evidence WHERE evidence_id = ?",
            (outcome.evidence_id,),
        ).fetchone()[0]
        assert str(tmp_path) not in evidence_payload
        assert "local_path" not in evidence_payload
    after_stat = pptx.stat()
    assert stream_sha256(pptx) == before_digest
    assert after_stat.st_size == before_stat.st_size
    assert after_stat.st_mtime_ns == before_stat.st_mtime_ns


def test_duplicate_materialization_reuses_exact_artifact_and_metadata(
    tmp_path: Path,
) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "duplicate.pptx", (png_bytes((1, 2, 3)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    store = ArtifactStore(tmp_path / ".artifacts")
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        first = materialize_source_visuals(
            catalog, registry, store,
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[0].version_id,),
            clock=lambda: NOW,
        )[0]
        second = materialize_source_visuals(
            catalog, registry, store,
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[0].version_id,),
            clock=lambda: NOW,
        )[0]
        assert first.status == second.status == "materialized"
        assert first.artifact_id == second.artifact_id
        assert first.reused is False
        assert second.reused is True
        assert catalog.connection.execute(
            "SELECT count(*) FROM artifact_metadata"
        ).fetchone()[0] == 1
        assert catalog.connection.execute(
            "SELECT count(*) FROM source_visual_artifacts"
        ).fetchone()[0] == 1


def test_identical_bytes_across_visuals_reuse_catalog_metadata_across_clocks(
    tmp_path: Path,
) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    payload = png_bytes((3, 4, 5))
    pptx = tmp_path / "same-bytes.pptx"
    presentation = Presentation()
    for index in range(2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(5), Inches(0.5)
        )
        box.text = f"Same bytes {index + 1}"
        slide.shapes.add_picture(
            BytesIO(payload), Inches(0.5), Inches(1), Inches(1), Inches(0.75)
        )
        slide.notes_slide.notes_text_frame.text = "Use the same source image."
    presentation.save(pptx)
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    assert len(result.visuals) == 2
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        first = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[0].version_id,),
            clock=lambda: NOW,
        )[0]
        second = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[1].version_id,),
            clock=lambda: datetime(2026, 7, 19, 12, 0, tzinfo=timezone.utc),
        )[0]
        assert first.artifact_id == second.artifact_id
        assert second.reused is True
        assert catalog.connection.execute(
            "SELECT count(*) FROM artifact_metadata"
        ).fetchone()[0] == 1
        assert catalog.connection.execute(
            "SELECT count(*) FROM source_visual_artifacts"
        ).fetchone()[0] == 2


def test_source_digest_change_fails_before_any_artifact_write(tmp_path: Path) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "changed.pptx", (png_bytes((5, 6, 7)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        with pptx.open("ab") as target:
            target.write(b"changed")
        outcomes = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[0].version_id,),
            clock=lambda: NOW,
        )
        assert outcomes[0].status == "failed"
        assert outcomes[0].error_code == "SOURCE_DIGEST_MISMATCH"
        assert catalog.connection.execute(
            "SELECT count(*) FROM artifact_metadata"
        ).fetchone()[0] == 0


def test_missing_relationship_does_not_roll_back_valid_sibling(tmp_path: Path) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "siblings.pptx", (png_bytes((8, 9, 10)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    valid = result.visuals[0]
    assert valid.source_locator is not None
    missing = valid.model_copy(
        update={
            "logical_id": "visual-missing-relationship",
            "version_id": "visual-missing-relationship-v1",
            "source_locator": valid.source_locator.model_copy(
                update={"relationship_id": "rId999"}
            ),
        }
    )
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        catalog.insert_visual(missing)
        outcomes = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(missing.version_id, valid.version_id),
            clock=lambda: NOW,
        )
        by_id = {outcome.visual_version_id: outcome for outcome in outcomes}
        assert by_id[missing.version_id].status == "failed"
        assert by_id[missing.version_id].error_code == "RELATIONSHIP_NOT_FOUND"
        assert by_id[valid.version_id].status == "materialized"
        assert catalog.connection.execute(
            "SELECT count(*) FROM artifact_metadata"
        ).fetchone()[0] == 1


def test_visual_digest_mismatch_is_sanitized_and_not_persisted(tmp_path: Path) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "visual-mismatch.pptx", (png_bytes((11, 12, 13)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    visual = result.visuals[0]
    forged = visual.model_copy(
        update={
            "logical_id": "visual-forged",
            "version_id": "visual-forged-v1",
            "content_digest": "f" * 64,
        }
    )
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        catalog.insert_source(result.source)
        catalog.insert_visual(forged)
        outcomes = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(forged.version_id,),
            clock=lambda: NOW,
        )
        assert outcomes[0].status == "failed"
        assert outcomes[0].error_code == "VISUAL_DIGEST_MISMATCH"
        assert str(tmp_path) not in (outcomes[0].message or "")
        assert catalog.connection.execute(
            "SELECT count(*) FROM source_visual_artifacts"
        ).fetchone()[0] == 0


def test_forged_visual_identity_cannot_claim_real_relationship_bytes(
    tmp_path: Path,
) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "forged-identity.pptx", (png_bytes((23, 24, 25)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    valid = result.visuals[0]
    forged = valid.model_copy(
        update={
            "logical_id": "forged-visual-logical",
            "version_id": "forged-visual-v1",
        }
    )
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        catalog.insert_visual(forged)
        outcome = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(forged.version_id,),
            clock=lambda: NOW,
        )[0]
        assert outcome.status == "failed"
        assert outcome.error_code == "VISUAL_IDENTITY_INVALID"
        assert catalog.connection.execute(
            "SELECT count(*) FROM source_visual_artifacts"
        ).fetchone()[0] == 0


def test_parser_pins_source_and_chunk_parents_for_materialization(tmp_path: Path) -> None:
    pptx = build_pptx(tmp_path / "parents.pptx", (png_bytes((14, 15, 16)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    visual = result.visuals[0]
    assert result.source.version_id in visual.derived_from_version_ids
    assert result.chunks[0].chunk_id in visual.derived_from_version_ids
    assert visual.source_locator is not None
    assert visual.source_locator.relationship_id
    assert visual.content_digest == hashlib.sha256(png_bytes((14, 15, 16))).hexdigest()


def test_duplicate_pptx_media_member_fails_closed(tmp_path: Path) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "duplicate-member.pptx", (png_bytes((17, 18, 19)),))
    with ZipFile(pptx, "a", ZIP_DEFLATED) as archive:
        media_name = next(
            info.filename for info in archive.infolist() if info.filename.startswith("ppt/media/")
        )
        payload = archive.read(media_name)
        archive.writestr(media_name, payload)
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        outcome = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[0].version_id,),
            clock=lambda: NOW,
        )[0]
        assert outcome.status == "failed"
        assert outcome.error_code == "RELATIONSHIP_INVALID"
        assert catalog.connection.execute(
            "SELECT count(*) FROM artifact_metadata"
        ).fetchone()[0] == 0


def test_artifact_and_source_visual_rows_are_foreign_key_bound_and_immutable(
    tmp_path: Path,
) -> None:
    from course_helper.artifacts import ArtifactStore
    from course_helper.source_visuals import materialize_source_visuals

    pptx = build_pptx(tmp_path / "immutable.pptx", (png_bytes((20, 21, 22)),))
    registry = SourceRootRegistry({"fixture": tmp_path})
    result = PptxParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path=pptx.name)
    )
    with KnowledgeCatalog.open(tmp_path / "catalog.db") as catalog:
        persist_extraction(catalog, result)
        outcome = materialize_source_visuals(
            catalog, registry, ArtifactStore(tmp_path / ".artifacts"),
            source_version_id=result.source.version_id,
            visual_version_ids=(result.visuals[0].version_id,),
            clock=lambda: NOW,
        )[0]
        assert outcome.artifact_id is not None
        assert catalog.connection.execute(
            "SELECT count(*) FROM lineage WHERE from_version_id = ? "
            "AND to_version_id = ? AND relation = 'derived_from'",
            (outcome.artifact_id, result.visuals[0].version_id),
        ).fetchone()[0] == 1
        with pytest.raises(sqlite3.IntegrityError, match="immutable artifact"):
            catalog.connection.execute(
                "UPDATE artifact_metadata SET byte_size = byte_size + 1 "
                "WHERE artifact_id = ?",
                (outcome.artifact_id,),
            )
        with pytest.raises(sqlite3.IntegrityError, match="immutable source visual"):
            catalog.connection.execute(
                "DELETE FROM source_visual_artifacts WHERE visual_version_id = ?",
                (result.visuals[0].version_id,),
            )
        with pytest.raises(sqlite3.IntegrityError):
            catalog.connection.execute(
                "INSERT INTO source_visual_artifacts("
                "materialization_id, visual_version_id, artifact_id, source_version_id, "
                "source_content_digest, visual_content_digest, slide_number, "
                "relationship_id, evidence_id, content_digest, payload_json, created_at"
                ") VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
                (
                    "bad-fk-materialization",
                    result.visuals[0].version_id,
                    "artifact-" + "f" * 64,
                    result.source.version_id,
                    result.source.content_digest,
                    result.visuals[0].content_digest,
                    1,
                    "rId1",
                    result.evidence.evidence_id,
                    "f" * 64,
                    "{}",
                    NOW.isoformat(),
                ),
            )
