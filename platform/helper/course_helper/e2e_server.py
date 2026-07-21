"""Isolated loopback Helper lifecycle used by the fixture-backed browser gate."""

from __future__ import annotations

import argparse
import json
import os
from datetime import datetime
from io import BytesIO
from pathlib import Path
from urllib.parse import urlsplit
from zipfile import ZipFile, ZipInfo

import uvicorn
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from course_helper.api import HelperRuntime, create_app
from course_helper.catalog import KnowledgeCatalog
from course_helper.jobs import BoundedJobRunner, WorkerRuntimeConfig
from course_helper.personal_supervisor import PersonalCourseSupervisor
from course_helper.session import LaunchSession


_CANONICAL_ZIP_TIMESTAMP = (2000, 1, 1, 0, 0, 0)


def _canonicalize_zip_archive(path: Path) -> None:
    canonical = path.with_suffix(f"{path.suffix}.canonical")
    with ZipFile(path, "r") as source, ZipFile(canonical, "w") as target:
        for member in sorted(source.infolist(), key=lambda item: item.filename):
            info = ZipInfo(member.filename, date_time=_CANONICAL_ZIP_TIMESTAMP)
            info.compress_type = member.compress_type
            info.create_system = member.create_system
            info.external_attr = member.external_attr
            target.writestr(info, source.read(member.filename))
    canonical.replace(path)


def _safe_origin(value: str) -> str:
    parsed = urlsplit(value)
    if (
        parsed.scheme != "http"
        or parsed.hostname != "127.0.0.1"
        or parsed.port is None
        or parsed.path
        or parsed.query
        or parsed.fragment
        or parsed.username is not None
        or parsed.password is not None
    ):
        raise ValueError("E2E web origin must be exact loopback HTTP")
    return value


def _write_fixtures(root: Path) -> dict[str, str]:
    fixtures = root / "fixtures"
    fixtures.mkdir(parents=True, exist_ok=True)
    markdown = fixtures / "evidence-course.md"
    markdown.write_text(
        "# 证据优先的 AI 课程\n\n## 学习目标\n\n理解真实来源、知识卡与课程发布的证据链。\n\n## 方法\n\n每个结论都需要引用、标签与固定版本。\n",
        encoding="utf-8",
    )
    dataset = fixtures / "segments.csv"
    dataset.write_text(
        "segment,revenue\n新客户,12\n成长客户,28\n核心客户,46\n",
        encoding="utf-8",
    )
    image = Image.new("RGB", (160, 100), (32, 101, 255))
    image_bytes = BytesIO()
    image.save(image_bytes, format="PNG")
    pptx = fixtures / "ai-evidence.pptx"
    presentation = Presentation()
    fixed_time = datetime(2000, 1, 1)
    presentation.core_properties.created = fixed_time
    presentation.core_properties.modified = fixed_time
    presentation.core_properties.author = "Course Studio E2E"
    presentation.core_properties.last_modified_by = "Course Studio E2E"
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(6), Inches(0.8)).text = "真实来源图形"
    slide.shapes.add_picture(BytesIO(image_bytes.getvalue()), Inches(1), Inches(1.5), width=Inches(3.2))
    presentation.save(pptx)
    _canonicalize_zip_archive(pptx)
    return {"markdown": str(markdown), "dataset": str(dataset), "pptx": str(pptx)}


def main() -> int:
    parser = argparse.ArgumentParser(prog="course-helper-e2e")
    parser.add_argument("--runtime-root", type=Path, required=True)
    parser.add_argument("--web-origin", required=True)
    parser.add_argument("--port", type=int, required=True)
    parser.add_argument("--launch-file", type=Path, required=True)
    args = parser.parse_args()
    if os.environ.get("COURSE_E2E_FIXTURE") != "1":
        raise SystemExit("E2E fixture is not authorized")
    runtime_root = args.runtime_root.resolve()
    launch_file = args.launch_file.resolve()
    if launch_file.parent != runtime_root or not 1 <= args.port <= 65535:
        raise SystemExit("E2E runtime paths are invalid")
    runtime_root.mkdir(parents=True, exist_ok=True)
    web_origin = _safe_origin(args.web_origin)
    fixture_paths = _write_fixtures(runtime_root)
    app_data = runtime_root / "app-data"
    database = runtime_root / "knowledge.db"
    app_data.mkdir(parents=True, exist_ok=True)
    with KnowledgeCatalog.open(database):
        pass
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(app_data),
        source_roots=(("e2e-fixtures", str(runtime_root / "fixtures")),),
        network_fixture=True,
    )
    session = LaunchSession.create(allowed_origin=web_origin, ttl_seconds=600)
    helper_origin = f"http://127.0.0.1:{args.port}"
    payload = {
        "schemaVersion": 1,
        "launchUrl": session.connect_url(
            web_application_url=web_origin,
            helper_base_url=helper_origin,
        ),
        "helperOrigin": helper_origin,
        "webOrigin": web_origin,
        "fixtures": fixture_paths,
    }
    temporary = launch_file.with_suffix(".tmp")
    temporary.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    temporary.replace(launch_file)
    supervisor = PersonalCourseSupervisor(config)
    supervisor.resume_pending()
    try:
        uvicorn.run(
            create_app(
                HelperRuntime(
                    config=config,
                    launch_session=session,
                    job_runner=BoundedJobRunner(config),
                    personal_course_supervisor=supervisor,
                )
            ),
            host="127.0.0.1",
            port=args.port,
            log_level="warning",
        )
    finally:
        supervisor.shutdown()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
