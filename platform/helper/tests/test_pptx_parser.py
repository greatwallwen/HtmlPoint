from __future__ import annotations

import hashlib
import os
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from course_helper.domain.common import SourceLocator
from course_helper.parsers.pptx_parser import PptxParser
from course_helper.source_roots import SourceRootRegistry


RELATIONSHIP_NAMESPACE = "http://schemas.openxmlformats.org/package/2006/relationships"


def _png_1x1() -> bytes:
    stream = BytesIO()
    Image.new("RGB", (1, 1), color=(36, 99, 235)).save(stream, format="PNG")
    return stream.getvalue()


PNG_1X1 = _png_1x1()


def build_small_pptx(
    path: Path,
    *,
    title: str | None,
    image_bytes: bytes,
    notes_text: str | None = "Presenter notes come first.",
) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    if title is not None:
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
        text_box.text = title
    slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.5), Inches(1.75), Inches(1), Inches(1))
    if notes_text is not None:
        slide.notes_slide.notes_text_frame.text = notes_text
    presentation.save(path)
    return path


def add_malformed_external_image_relationship(path: Path) -> str:
    relationship_id = "rId999"
    relationship_name = "ppt/slides/_rels/slide1.xml.rels"
    rewritten = path.with_name(f"{path.stem}-rewritten.pptx")
    with ZipFile(path, "r") as source, ZipFile(rewritten, "w", ZIP_DEFLATED) as target:
        for info in source.infolist():
            payload = source.read(info.filename)
            if info.filename == relationship_name:
                root = ElementTree.fromstring(payload)
                ElementTree.SubElement(
                    root,
                    f"{{{RELATIONSHIP_NAMESPACE}}}Relationship",
                    {
                        "Id": relationship_id,
                        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
                        "Target": "https://example.invalid/missing.png",
                        "TargetMode": "External",
                    },
                )
                payload = ElementTree.tostring(root, encoding="utf-8", xml_declaration=True)
            target.writestr(info, payload)
    rewritten.replace(path)
    return relationship_id


def parser_for(root: Path) -> PptxParser:
    return PptxParser(SourceRootRegistry({"fixture": root}))


def demo_parser() -> PptxParser:
    configured_root = os.environ.get("COURSE_REFERENCE_ROOT")
    if not configured_root:
        pytest.skip("COURSE_REFERENCE_ROOT is required for reference_demo tests")
    reference_root = Path(configured_root)
    return PptxParser(SourceRootRegistry({"reference-demo": reference_root}))


def test_pptx_parser_keeps_notes_slide_text_and_visual_relationship(tmp_path: Path) -> None:
    fixture = build_small_pptx(
        tmp_path / "fixture.pptx",
        title="Transformer",
        image_bytes=PNG_1X1,
    )

    result = parser_for(tmp_path).parse(
        SourceLocator(root_id="fixture", relative_path=fixture.name)
    )

    assert [chunk.locator.slide_number for chunk in result.chunks] == [1]
    assert result.chunks[0].notes_text == "Presenter notes come first."
    assert result.chunks[0].slide_text == "Transformer"
    assert result.chunks[0].normalized_text.startswith("Presenter notes come first.")
    assert result.chunks[0].normalized_text.endswith("Transformer")
    assert result.visuals[0].source_locator is not None
    assert result.visuals[0].source_locator.slide_number == 1
    assert result.visuals[0].source_locator.relationship_id
    assert result.visuals[0].content_digest == hashlib.sha256(PNG_1X1).hexdigest()
    assert (result.visuals[0].width, result.visuals[0].height) == (1, 1)
    assert result.chunks[0].media_version_ids == (result.visuals[0].version_id,)


def test_image_only_slide_emits_visual_without_empty_chunk(tmp_path: Path) -> None:
    fixture = build_small_pptx(
        tmp_path / "image-only.pptx",
        title=None,
        image_bytes=PNG_1X1,
        notes_text=None,
    )

    result = parser_for(tmp_path).parse(
        SourceLocator(root_id="fixture", relative_path=fixture.name)
    )

    assert result.chunks == ()
    assert len(result.visuals) == 1
    assert result.visuals[0].content_digest == hashlib.sha256(PNG_1X1).hexdigest()
    assert result.visuals[0].source_locator is not None
    assert result.visuals[0].source_locator.slide_number == 1
    assert any(
        check.code == "pptx-image-relationship" and check.status == "passed"
        for check in result.evidence.checks
    )


def test_image_only_slide_records_sanitized_malformed_relationship(
    tmp_path: Path,
) -> None:
    fixture = build_small_pptx(
        tmp_path / "image-only-malformed.pptx",
        title=None,
        image_bytes=PNG_1X1,
        notes_text=None,
    )
    malformed_relationship_id = add_malformed_external_image_relationship(fixture)

    result = parser_for(tmp_path).parse(
        SourceLocator(root_id="fixture", relative_path=fixture.name)
    )

    assert len(result.visuals) == 1
    failed_media_check = next(
        check
        for check in result.evidence.checks
        if check.code == "pptx-image-relationship" and check.status == "failed"
    )
    assert failed_media_check.details["relationship_id"] == malformed_relationship_id
    assert failed_media_check.details["exception_type"] == "ValueError"
    assert "example.invalid" not in repr(dict(failed_media_check.details))


def test_pptx_parser_is_deterministic_for_the_same_source(tmp_path: Path) -> None:
    fixture = build_small_pptx(
        tmp_path / "fixture.pptx",
        title="Deterministic extraction",
        image_bytes=PNG_1X1,
    )
    parser = parser_for(tmp_path)
    locator = SourceLocator(root_id="fixture", relative_path=fixture.name)

    assert parser.parse(locator) == parser.parse(locator)


def test_malformed_image_relationship_records_failure_without_losing_text(
    tmp_path: Path,
) -> None:
    fixture = build_small_pptx(
        tmp_path / "fixture.pptx",
        title="Text survives media failure",
        image_bytes=PNG_1X1,
    )
    malformed_relationship_id = add_malformed_external_image_relationship(fixture)

    result = parser_for(tmp_path).parse(
        SourceLocator(root_id="fixture", relative_path=fixture.name)
    )

    assert result.chunks[0].slide_text == "Text survives media failure"
    assert len(result.visuals) == 1
    assert result.source.extraction_status == "partial"
    assert result.evidence.status == "degraded"
    failed_checks = [check for check in result.evidence.checks if check.status == "failed"]
    assert len(failed_checks) == 1
    assert failed_checks[0].code == "pptx-image-relationship"
    assert failed_checks[0].details["slide_number"] == 1
    assert failed_checks[0].details["relationship_id"] == malformed_relationship_id


@pytest.mark.reference_demo
def test_ai_pptx_demo_extracts_notes_first() -> None:
    result = demo_parser().parse(
        SourceLocator(root_id="reference-demo", relative_path="AI.pptx"),
        range(3, 19),
    )

    assert len(result.chunks) == 16
    assert all(chunk.locator.slide_number in range(3, 19) for chunk in result.chunks)
    assert all(chunk.notes_text.strip() for chunk in result.chunks)
    assert any(
        visual.source_locator is not None and visual.source_locator.slide_number == 3
        for visual in result.visuals
    )
