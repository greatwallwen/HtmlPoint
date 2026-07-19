from __future__ import annotations

import hashlib
import json
import os
import sqlite3
from datetime import datetime, timezone
from pathlib import Path

import pytest
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


DEMO_SOURCE_PATHS = (
    "AI.pptx",
    "AIGC实操 -数据分析.md",
    "AIGC实操-Prompt工程.md",
    "dataset/1-train.csv",
    "AIExcelData/ex-17-RFM.xlsx",
)


def synthetic_source_state(root: Path) -> tuple[object, ...]:
    entries = tuple(
        sorted(
            (
                path.relative_to(root).as_posix(),
                path.is_dir(),
                path.stat().st_size,
                path.stat().st_mtime_ns,
            )
            for path in root.rglob("*")
        )
    )
    sources = tuple(
        (
            relative_path,
            (root / relative_path).read_bytes(),
            (root / relative_path).stat().st_size,
            (root / relative_path).stat().st_mtime_ns,
        )
        for relative_path in DEMO_SOURCE_PATHS
    )
    return entries, sources


def output_path_state(path: Path) -> tuple[object, ...]:
    if not os.path.lexists(path):
        return ("missing",)
    metadata = path.lstat()
    payload = path.read_bytes() if path.is_file() else None
    return (
        "existing",
        path.is_symlink(),
        metadata.st_size,
        metadata.st_mtime_ns,
        metadata.st_nlink,
        payload,
    )


def test_catalog_counts_use_lifecycle_projections_without_reference_access(
    tmp_path: Path,
) -> None:
    from course_helper.cards import VOCABULARY_VERSION_ID, seed_vocabulary
    from course_helper.catalog import KnowledgeCatalog
    from course_helper.demo import _catalog_counts
    from course_helper.domain.common import ActorRef
    from course_helper.domain.knowledge import CardContentNode, KnowledgeCardVersion
    from course_helper.lifecycle import append_card_lifecycle_event

    now = datetime(2026, 7, 17, tzinfo=timezone.utc)
    with KnowledgeCatalog.open(tmp_path / "projection-counts.sqlite3") as catalog:
        seed_vocabulary(catalog)
        published = catalog.insert_card(
            KnowledgeCardVersion(
                logical_id="demo-count-card",
                version_id="demo-count-card-v1",
                revision=1,
                content_digest=hashlib.sha256(b"demo-count-card").hexdigest(),
                created_at=now,
                created_by=ActorRef(actor_type="system", actor_id="demo-count-tests"),
                main_type_id="concept",
                title="Projected count card",
                learning_objective="Count only effective lifecycle states",
                content_ast=(CardContentNode(type="paragraph", text="Projection truth"),),
                suggested_minutes=5,
                vocabulary_version_id=VOCABULARY_VERSION_ID,
                status="review",
            )
        )
        append_card_lifecycle_event(
            catalog.connection,
            card_version_id=published.version_id,
            event_id="publish-demo-count-card",
            request_digest="0" * 64,
            event_type="publish",
            occurred_at=now,
            actor_id="demo-count-tests",
        )
        assert _catalog_counts(catalog)["published_cards"] == 1
        append_card_lifecycle_event(
            catalog.connection,
            card_version_id=published.version_id,
            event_id="suspend-demo-count-card",
            request_digest="4" * 64,
            event_type="suspend",
            occurred_at=now,
            actor_id="demo-count-tests",
        )
        assert _catalog_counts(catalog)["published_cards"] == 0
        append_card_lifecycle_event(
            catalog.connection,
            card_version_id=published.version_id,
            event_id="archive-demo-count-card",
            request_digest="5" * 64,
            event_type="archive",
            occurred_at=now,
            actor_id="demo-count-tests",
        )
        counts = _catalog_counts(catalog)
        assert counts["published_cards"] == 0
        assert counts["archived_cards"] == 1


def assert_invalid_demo_outputs_leave_everything_untouched(
    monkeypatch,
    *,
    source_root: Path,
    database_path: Path,
    evidence_path: Path,
) -> None:
    from course_helper.demo import DemoConfigurationError, run_reference_demo

    source_before = synthetic_source_state(source_root)
    output_before = {
        path: output_path_state(path) for path in {database_path, evidence_path}
    }
    source_opens: list[str] = []
    original_open = Path.open

    def recording_open(path: Path, *args, **kwargs):
        resolved = path.resolve()
        if resolved.is_relative_to(source_root.resolve()):
            source_opens.append(resolved.relative_to(source_root.resolve()).as_posix())
        return original_open(path, *args, **kwargs)

    monkeypatch.setattr(Path, "open", recording_open)
    try:
        with pytest.raises(
            DemoConfigurationError,
            match="unsafe-demo-output",
        ) as captured_error:
            run_reference_demo(
                source_root,
                database_path,
                evidence_path,
            )
    finally:
        monkeypatch.setattr(Path, "open", original_open)

    assert source_opens == []
    error_text = str(captured_error.value)
    assert str(source_root) not in error_text
    assert str(database_path) not in error_text
    assert str(evidence_path) not in error_text
    assert synthetic_source_state(source_root) == source_before
    assert {
        path: output_path_state(path) for path in {database_path, evidence_path}
    } == output_before


def build_synthetic_reference_root(root: Path) -> Path:
    (root / "dataset").mkdir(parents=True)
    (root / "AIExcelData" / "weights").mkdir(parents=True)
    (root / "AIExcelData" / "assets").mkdir(parents=True)

    presentation = Presentation()
    for slide_number in range(1, 19):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            Inches(5),
            Inches(1),
        )
        text_box.text = "人工智能基础"
        slide.notes_slide.notes_text_frame.text = (
            f"人工智能课程证据，第 {slide_number} 张。"
        )
    presentation.save(root / "AI.pptx")

    (root / "AIGC实操 -数据分析.md").write_text(
        "# 自行车共享需求\n自行车共享需求分析使用天气与租赁数据。\n"
        "## 分析步骤\n这是选中单元内保留的后代分块。\n",
        encoding="utf-8",
    )
    (root / "AIGC实操-Prompt工程.md").write_text(
        "# Prompt概论\nPrompt用于描述任务。\n"
        "## 正确提问\n正确提问需要明确目标。\n"
        "![non-allowlisted](AIExcelData/assets/unlisted.png)\n",
        encoding="utf-8",
    )
    (root / "dataset" / "1-train.csv").write_text(
        "record_id,value\n1,10\n2,20\n",
        encoding="utf-8",
    )

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "RFM"
    worksheet.append(("customer_id", "recency", "frequency", "monetary"))
    worksheet.append((1, 5, 2, 100.0))
    worksheet.append((2, 3, 4, 250.0))
    workbook.save(root / "AIExcelData" / "ex-17-RFM.xlsx")
    workbook.close()

    (root / "AIExcelData" / "weights" / "sam_vit_h_4b8939.pth").write_bytes(
        b"metadata-only fixture"
    )
    (root / "AIExcelData" / "assets" / "unlisted.png").write_bytes(
        bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c489"
            "0000000d49444154789c6360f8cff00000040101089d1d0d0000000049454e44"
            "ae426082"
        )
    )
    return root


def test_output_validation_rejects_database_inside_source_before_any_deep_read(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper.demo import DemoConfigurationError, run_reference_demo

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    before = synthetic_source_state(source_root)
    source_opens: list[str] = []
    original_open = Path.open

    def recording_open(path: Path, *args, **kwargs):
        resolved = path.resolve()
        if resolved.is_relative_to(source_root.resolve()):
            source_opens.append(resolved.relative_to(source_root.resolve()).as_posix())
        return original_open(path, *args, **kwargs)

    monkeypatch.setattr(Path, "open", recording_open)
    evidence_path = tmp_path / "new-output-parent" / "receipt.json"

    with pytest.raises(DemoConfigurationError, match="unsafe-demo-output"):
        run_reference_demo(
            source_root,
            source_root / "forbidden.db",
            evidence_path,
        )

    monkeypatch.setattr(Path, "open", original_open)
    assert source_opens == []
    assert synthetic_source_state(source_root) == before
    assert not evidence_path.parent.exists()
    assert not (source_root / "forbidden.db").exists()


def test_output_validation_suppresses_path_bearing_resolution_causes(
    tmp_path: Path,
) -> None:
    from course_helper.demo import DemoConfigurationError, validate_demo_outputs

    missing_root = tmp_path / "missing-reference-root"
    with pytest.raises(
        DemoConfigurationError,
        match="unsafe-demo-output: path resolution failed",
    ) as captured_error:
        validate_demo_outputs(
            missing_root,
            tmp_path / "knowledge.db",
            tmp_path / "receipt.json",
        )

    assert str(tmp_path) not in str(captured_error.value)
    assert captured_error.value.__cause__ is None


def test_output_validation_rejects_evidence_inside_source_before_any_deep_read(
    monkeypatch,
    tmp_path: Path,
) -> None:
    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    assert_invalid_demo_outputs_leave_everything_untouched(
        monkeypatch,
        source_root=source_root,
        database_path=tmp_path / "knowledge.db",
        evidence_path=source_root / "forbidden-receipt.json",
    )


def test_output_validation_rejects_one_path_for_database_and_evidence(
    monkeypatch,
    tmp_path: Path,
) -> None:
    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    shared_path = tmp_path / "new-output-parent" / "shared-output"
    assert_invalid_demo_outputs_leave_everything_untouched(
        monkeypatch,
        source_root=source_root,
        database_path=shared_path,
        evidence_path=shared_path,
    )
    assert not shared_path.parent.exists()


@pytest.mark.parametrize("sidecar_suffix", ["-wal", "-shm", "-journal"])
def test_output_validation_rejects_evidence_at_a_database_sidecar(
    monkeypatch,
    tmp_path: Path,
    sidecar_suffix: str,
) -> None:
    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "knowledge.db"
    assert_invalid_demo_outputs_leave_everything_untouched(
        monkeypatch,
        source_root=source_root,
        database_path=database_path,
        evidence_path=Path(f"{database_path}{sidecar_suffix}"),
    )


def test_output_validation_rejects_database_hardlinked_to_an_allowlisted_source(
    monkeypatch,
    tmp_path: Path,
) -> None:
    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "source-hardlink.db"
    try:
        os.link(source_root / "AI.pptx", database_path)
    except OSError as error:
        pytest.skip(f"hardlinks unavailable: {type(error).__name__}")
    assert_invalid_demo_outputs_leave_everything_untouched(
        monkeypatch,
        source_root=source_root,
        database_path=database_path,
        evidence_path=tmp_path / "receipt.json",
    )


def test_output_validation_rejects_existing_hardlinked_outputs(
    monkeypatch,
    tmp_path: Path,
) -> None:
    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "shared-database.db"
    evidence_path = tmp_path / "shared-receipt.json"
    database_path.write_bytes(b"")
    try:
        os.link(database_path, evidence_path)
    except OSError as error:
        pytest.skip(f"hardlinks unavailable: {type(error).__name__}")
    assert_invalid_demo_outputs_leave_everything_untouched(
        monkeypatch,
        source_root=source_root,
        database_path=database_path,
        evidence_path=evidence_path,
    )


def test_output_validation_resolves_an_existing_parent_link(
    tmp_path: Path,
) -> None:
    from course_helper.demo import validate_demo_outputs

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    actual_parent = tmp_path / "actual-output"
    linked_parent = tmp_path / "linked-output"
    actual_parent.mkdir()
    try:
        linked_parent.symlink_to(actual_parent, target_is_directory=True)
    except OSError as error:
        import inspect

        source = inspect.getsource(validate_demo_outputs)
        assert ".resolve(" in source
        pytest.skip(f"directory links unavailable: {type(error).__name__}")

    outputs = validate_demo_outputs(
        source_root,
        linked_parent / "knowledge.db",
        linked_parent / "receipt.json",
    )

    assert outputs.database_path == (actual_parent / "knowledge.db").resolve()
    assert outputs.evidence_path == (actual_parent / "receipt.json").resolve()


def test_output_validation_rejects_an_existing_output_symlink(
    tmp_path: Path,
) -> None:
    from course_helper.demo import DemoConfigurationError, validate_demo_outputs

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    target_path = tmp_path / "real-database.db"
    linked_path = tmp_path / "linked-database.db"
    target_path.write_bytes(b"user output")
    try:
        linked_path.symlink_to(target_path)
    except OSError as error:
        import inspect

        source = inspect.getsource(validate_demo_outputs)
        assert "_resolve_demo_output" in source
        pytest.skip(f"file links unavailable: {type(error).__name__}")

    with pytest.raises(DemoConfigurationError, match="unsafe-demo-output"):
        validate_demo_outputs(
            source_root,
            linked_path,
            tmp_path / "receipt.json",
        )

    assert linked_path.is_symlink()
    assert target_path.read_bytes() == b"user output"


def test_output_validation_rejects_an_existing_database_sidecar(
    tmp_path: Path,
) -> None:
    from course_helper.catalog import KnowledgeCatalog
    from course_helper.demo import DemoConfigurationError, validate_demo_outputs

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database_path):
        pass
    sidecar_path = Path(f"{database_path}-wal")
    sidecar_path.write_bytes(b"pre-existing user sidecar")
    database_before = database_path.read_bytes()
    sidecar_before = sidecar_path.read_bytes()

    with pytest.raises(DemoConfigurationError, match="unsafe-demo-output"):
        validate_demo_outputs(
            source_root,
            database_path,
            tmp_path / "receipt.json",
        )

    assert database_path.read_bytes() == database_before
    assert sidecar_path.read_bytes() == sidecar_before


def test_source_change_after_parsing_does_not_mutate_the_target_database(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper import demo
    from course_helper.catalog import KnowledgeCatalog
    from course_helper.demo import DemoIntegrityError, run_reference_demo

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "knowledge.db"
    evidence_path = tmp_path / "failed-receipt.json"
    with KnowledgeCatalog.open(database_path):
        pass
    database_before = database_path.read_bytes()
    with sqlite3.connect(database_path) as connection:
        row_counts_before = tuple(
            connection.execute(f"SELECT count(*) FROM {table}").fetchone()[0]
            for table in ("sources", "chunks", "datasets", "cards", "evidence")
        )
    original_parse = demo._parse_allowlisted_sources

    def parse_then_change_source(*args, **kwargs):
        result = original_parse(*args, **kwargs)
        prompt_path = source_root / "AIGC实操-Prompt工程.md"
        prompt_path.write_text(
            prompt_path.read_text(encoding="utf-8") + "\n源在解析后发生变化。\n",
            encoding="utf-8",
        )
        return result

    monkeypatch.setattr(demo, "_parse_allowlisted_sources", parse_then_change_source)

    with pytest.raises(DemoIntegrityError, match="forbidden-source-write"):
        run_reference_demo(source_root, database_path, evidence_path)

    assert database_path.read_bytes() == database_before
    with sqlite3.connect(database_path) as connection:
        assert tuple(
            connection.execute(f"SELECT count(*) FROM {table}").fetchone()[0]
            for table in ("sources", "chunks", "datasets", "cards", "evidence")
        ) == row_counts_before
    failed_receipt = json.loads(evidence_path.read_text(encoding="utf-8"))
    assert failed_receipt["status"] == "failed"
    assert failed_receipt["forbidden_source_writes"] > 0
    assert str(source_root) not in evidence_path.read_text(encoding="utf-8")
    assert list(tmp_path.glob(".*demo-stage*")) == []


def test_receipt_atomic_replace_failure_raises_without_success_or_temp_file(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper import demo
    from course_helper.demo import run_reference_demo

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "knowledge.db"
    evidence_path = tmp_path / "receipt.json"
    original_replace = demo.os.replace

    def fail_receipt_replace(source, destination):
        if Path(destination) == evidence_path.resolve():
            raise OSError("injected receipt replace failure")
        return original_replace(source, destination)

    monkeypatch.setattr(demo.os, "replace", fail_receipt_replace)

    with pytest.raises(OSError, match="injected receipt replace failure"):
        run_reference_demo(source_root, database_path, evidence_path)

    assert database_path.exists()
    assert not evidence_path.exists()
    assert list(tmp_path.glob(".*receipt*.tmp")) == []
    assert list(tmp_path.glob(".*demo-stage*")) == []


def test_database_promotion_failure_preserves_target_and_cleans_stage(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper import demo
    from course_helper.catalog import KnowledgeCatalog
    from course_helper.demo import run_reference_demo

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "knowledge.db"
    evidence_path = tmp_path / "receipt.json"
    with KnowledgeCatalog.open(database_path):
        pass
    database_before = database_path.read_bytes()
    original_replace = demo.os.replace

    def fail_database_promotion(source, destination):
        if (
            Path(destination) == database_path.resolve()
            and ".demo-stage-" in Path(source).name
        ):
            raise OSError("injected database promotion failure")
        return original_replace(source, destination)

    monkeypatch.setattr(demo.os, "replace", fail_database_promotion)

    with pytest.raises(OSError, match="injected database promotion failure"):
        run_reference_demo(source_root, database_path, evidence_path)

    assert database_path.read_bytes() == database_before
    assert not evidence_path.exists()
    assert list(tmp_path.glob(".*demo-stage*")) == []


def test_stage_write_exception_preserves_target_rows_and_cleans_stage(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper import demo
    from course_helper.catalog import KnowledgeCatalog
    from course_helper.demo import run_reference_demo

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "knowledge.db"
    evidence_path = tmp_path / "receipt.json"
    with KnowledgeCatalog.open(database_path):
        pass
    database_before = database_path.read_bytes()
    original_persist = demo._persist_parsed_objects

    def persist_then_fail(*args, **kwargs):
        original_persist(*args, **kwargs)
        raise RuntimeError("injected staged write failure")

    monkeypatch.setattr(demo, "_persist_parsed_objects", persist_then_fail)

    with pytest.raises(RuntimeError, match="injected staged write failure"):
        run_reference_demo(source_root, database_path, evidence_path)

    assert database_path.read_bytes() == database_before
    with sqlite3.connect(database_path) as connection:
        assert tuple(
            connection.execute(f"SELECT count(*) FROM {table}").fetchone()[0]
            for table in ("sources", "chunks", "datasets", "cards", "evidence")
        ) == (0, 0, 0, 0, 0)
    assert not evidence_path.exists()
    assert list(tmp_path.glob(".*demo-stage*")) == []


def test_reference_demo_manifest_is_the_exact_utf8_allowlist() -> None:
    from course_helper.demo import load_demo_manifest

    manifest = load_demo_manifest()

    assert manifest.schema_version == 1
    assert manifest.root_id == "reference-demo"
    assert manifest.inventory_roots == ("dataset", "AIExcelData")
    assert tuple((source.kind, source.path) for source in manifest.sources) == (
        ("pptx", "AI.pptx"),
        ("markdown", "AIGC实操 -数据分析.md"),
        ("markdown", "AIGC实操-Prompt工程.md"),
        ("csv", "dataset/1-train.csv"),
        ("xlsx", "AIExcelData/ex-17-RFM.xlsx"),
    )
    assert manifest.sources[0].slides is not None
    assert manifest.sources[0].slides.start == 3
    assert manifest.sources[0].slides.end_inclusive == 18
    assert manifest.sources[1].headings == ("自行车共享需求",)
    assert manifest.sources[2].headings == ("Prompt概论", "正确提问")
    assert manifest.quarantined_extensions == (".pth", ".pt", ".tmp", ".whl")


def test_integrity_capture_hashes_only_the_allowlist_and_inventories_by_metadata(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper.demo import capture_demo_integrity

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    resolved_root = source_root.resolve()
    source_opens: list[str] = []
    original_open = Path.open

    def recording_open(path: Path, *args, **kwargs):
        resolved = path.resolve()
        if resolved.is_relative_to(resolved_root):
            source_opens.append(resolved.relative_to(resolved_root).as_posix())
        return original_open(path, *args, **kwargs)

    monkeypatch.setattr(Path, "open", recording_open)

    snapshot = capture_demo_integrity(source_root)

    assert source_opens == [
        "AI.pptx",
        "AIGC实操 -数据分析.md",
        "AIGC实操-Prompt工程.md",
        "dataset/1-train.csv",
        "AIExcelData/ex-17-RFM.xlsx",
    ]
    assert snapshot.root_id == "reference-demo"
    assert snapshot.inventory_integrity_scope == "metadata-only"
    assert snapshot.inventory_root_count == 2
    assert snapshot.quarantined_extension_counts[".pth"] == 1
    assert all(item.sha256 for item in snapshot.sources)
    assert str(source_root) not in snapshot.model_dump_json()


def test_synthetic_demo_orchestrates_traceable_read_only_knowledge(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper.demo import run_reference_demo

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    unlisted_image = (source_root / "AIExcelData" / "assets" / "unlisted.png").resolve()
    source_opens: list[Path] = []
    original_open = Path.open

    def recording_open(path: Path, *args, **kwargs):
        resolved = path.resolve()
        if resolved.is_relative_to(source_root.resolve()):
            source_opens.append(resolved)
        return original_open(path, *args, **kwargs)

    monkeypatch.setattr(Path, "open", recording_open)
    database_path = tmp_path / "knowledge.db"
    evidence_path = tmp_path / "receipt.json"

    receipt = run_reference_demo(source_root, database_path, evidence_path)

    assert receipt.root_id == "reference-demo"
    assert receipt.deep_read_source_count == 5
    assert receipt.hash_verified_source_count == 5
    assert receipt.inventory_root_count == 2
    assert receipt.inventory_integrity_scope == "metadata-only"
    assert receipt.pptx_slide_chunks == 16
    assert receipt.pptx_chunks_with_notes == 16
    assert receipt.markdown_units == {
        "自行车共享需求",
        "Prompt概论",
        "正确提问",
    }
    assert receipt.profiled_datasets == {
        "dataset/1-train.csv",
        "AIExcelData/ex-17-RFM.xlsx",
    }
    assert receipt.quarantined_extension_counts[".pth"] == 1
    assert receipt.published_card_count > 0
    assert receipt.review_decision_count > 0
    assert receipt.new_source_versions == 5
    assert receipt.new_card_count > 0
    assert receipt.forbidden_source_writes == 0
    assert all(item.changed_item_count == 0 for item in receipt.inventory_integrity)
    assert unlisted_image not in source_opens
    assert all(
        item.before_sha256 == item.after_sha256
        and item.before_metadata == item.after_metadata
        for item in receipt.source_integrity
    )
    assert len(receipt.retrievals) == 3
    assert all(item.hit_count > 0 for item in receipt.retrievals)
    assert all(item.evidence_status == "degraded" for item in receipt.retrievals)
    assert receipt.status == "degraded"

    receipt_text = evidence_path.read_text(encoding="utf-8")
    assert str(source_root) not in receipt_text
    assert "demo-stage" not in receipt_text
    assert json.loads(receipt_text)["root_id"] == "reference-demo"
    with sqlite3.connect(database_path) as connection:
        assert connection.execute(
            "SELECT count(*) FROM evidence WHERE kind = 'validation'"
        ).fetchone()[0] == receipt.review_decision_count
        assert connection.execute(
            "SELECT count(*) FROM review_tasks WHERE status = 'open'"
        ).fetchone()[0] == 0


def test_synthetic_demo_second_pass_has_no_new_or_duplicate_objects(
    tmp_path: Path,
) -> None:
    from course_helper.demo import run_reference_demo

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    database_path = tmp_path / "knowledge.db"
    first = run_reference_demo(
        source_root,
        database_path,
        tmp_path / "receipt-1.json",
    )
    second = run_reference_demo(
        source_root,
        database_path,
        tmp_path / "receipt-2.json",
    )

    assert first.new_source_versions == 5
    assert first.new_card_count > 0
    assert second.new_source_versions == 0
    assert second.new_card_count == 0
    assert second.new_evidence_count == 0
    assert second.duplicate_card_count == 0
    assert second.forbidden_source_writes == 0
    assert second.object_digests == first.object_digests


def test_cli_verify_idempotence_writes_canonical_second_pass_block(
    tmp_path: Path,
) -> None:
    from course_helper.demo import main

    source_root = build_synthetic_reference_root(tmp_path / "reference-root")
    evidence_path = tmp_path / "receipt.json"

    exit_code = main(
        [
            "--source-root",
            str(source_root),
            "--database",
            str(tmp_path / "knowledge.db"),
            "--evidence",
            str(evidence_path),
            "--verify-idempotence",
        ]
    )

    assert exit_code == 0
    receipt = json.loads(evidence_path.read_text(encoding="utf-8"))
    assert receipt["new_source_versions"] == 0
    assert receipt["new_card_count"] == 0
    assert receipt["idempotence"]["verified"] is True
    assert receipt["idempotence"]["pass_count"] == 2
    assert receipt["idempotence"]["second_pass"] == {
        "duplicate_card_count": 0,
        "forbidden_source_writes": 0,
        "new_card_count": 0,
        "new_evidence_count": 0,
        "new_source_versions": 0,
    }
    assert evidence_path.read_bytes().endswith(b"\n")


@pytest.mark.reference_demo
def test_reference_demo_builds_traceable_read_only_knowledge(
    monkeypatch,
    tmp_path: Path,
) -> None:
    from course_helper.demo import run_reference_demo

    configured_root = os.environ.get("COURSE_REFERENCE_ROOT")
    if not configured_root:
        pytest.skip("COURSE_REFERENCE_ROOT is required for the reference fixture")
    source_root = Path(configured_root).resolve()
    allowed_paths = {
        "AI.pptx",
        "AIGC实操 -数据分析.md",
        "AIGC实操-Prompt工程.md",
        "dataset/1-train.csv",
        "AIExcelData/ex-17-RFM.xlsx",
    }
    opened_paths: list[str] = []
    original_open = Path.open

    def recording_open(path: Path, *args, **kwargs):
        resolved = path.resolve()
        if resolved.is_relative_to(source_root):
            opened_paths.append(resolved.relative_to(source_root).as_posix())
        return original_open(path, *args, **kwargs)

    monkeypatch.setattr(Path, "open", recording_open)
    receipt = run_reference_demo(
        source_root,
        tmp_path / "knowledge.db",
        tmp_path / "receipt.json",
    )

    assert receipt.pptx_slide_chunks == 16
    assert receipt.pptx_chunks_with_notes == 16
    assert receipt.markdown_units == {
        "自行车共享需求",
        "Prompt概论",
        "正确提问",
    }
    assert receipt.profiled_datasets == {
        "dataset/1-train.csv",
        "AIExcelData/ex-17-RFM.xlsx",
    }
    assert receipt.quarantined_extension_counts[".pth"] >= 1
    assert receipt.published_card_count > 0
    assert receipt.forbidden_source_writes == 0
    assert all(item.changed_item_count == 0 for item in receipt.inventory_integrity)
    assert opened_paths
    assert set(opened_paths) <= allowed_paths
