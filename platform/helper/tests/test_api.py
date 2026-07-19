from __future__ import annotations

import asyncio
import hashlib
import json
import sqlite3
import threading
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime, timezone
from pathlib import Path

import pytest
import duckdb
from fastapi.testclient import TestClient
from pydantic import TypeAdapter, ValidationError

from course_helper.catalog import KnowledgeCatalog
from course_helper.domain.evidence import EvidenceCheck, EvidenceObject


def test_job_specs_use_discriminated_lower_camel_http_models() -> None:
    from course_helper.jobs import JobSpec, SourceIngestJob

    parsed = TypeAdapter(JobSpec).validate_python(
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "demo.pptx"},
            "selection": {"slideNumbers": [3, 4]},
        }
    )

    assert isinstance(parsed, SourceIngestJob)
    assert parsed.locator.root_id == "fixture"
    assert parsed.selection.slide_numbers == (3, 4)
    assert parsed.model_dump(mode="json", by_alias=True)["locator"] == {
        "rootId": "fixture",
        "relativePath": "demo.pptx",
    }
    with pytest.raises(ValidationError):
        TypeAdapter(JobSpec).validate_python(
            {
                "type": "source_ingest",
                "locator": {"root_id": "fixture", "relative_path": "demo.pptx"},
                "selection": {"slideNumbers": [3, 4]},
            }
        )


@pytest.mark.parametrize("relative_path", ("../secret.md", "C:/secret.md", "/secret.md"))
def test_http_source_locator_rejects_root_escape_before_dispatch(
    relative_path: str,
) -> None:
    from course_helper.jobs import JobSpec

    with pytest.raises(ValidationError):
        TypeAdapter(JobSpec).validate_python(
            {
                "type": "source_ingest",
                "locator": {"rootId": "fixture", "relativePath": relative_path},
                "selection": {"headingSelectors": []},
            }
        )


def test_all_job_boundaries_use_nested_lower_camel_and_forbid_extras() -> None:
    from course_helper.jobs import JobSpec

    adapter = TypeAdapter(JobSpec)
    payloads = (
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "lesson.md"},
            "selection": {"slideNumbers": [], "headingSelectors": ["Lesson"]},
        },
        {
            "type": "dataset_profile",
            "locator": {"rootId": "fixture", "relativePath": "data.parquet"},
            "sampleLimit": 20,
            "sheetName": "Sheet 1",
        },
        {
            "type": "knowledge_retrieve",
            "query": "RFM",
            "requiredTagIds": ["topic:data-analysis"],
            "limit": 5,
        },
        {
            "type": "knowledge_publish",
            "cardVersionId": "card-version-1",
        },
    )

    for payload in payloads:
        parsed = adapter.validate_python(payload)
        dumped = parsed.model_dump(mode="json", by_alias=True)
        assert dumped["type"] == payload["type"]
        assert not any("_" in key for key in dumped)
    assert adapter.validate_python(payloads[0]).model_dump(
        mode="json", by_alias=True
    )["selection"] == {
        "slideNumbers": [],
        "headingSelectors": ["Lesson"],
    }
    assert adapter.validate_python(payloads[1]).model_dump(
        mode="json", by_alias=True
    )["locator"] == {"rootId": "fixture", "relativePath": "data.parquet"}

    for invalid in (
        {
            **payloads[0],
            "locator": {**payloads[0]["locator"], "command": "whoami"},
        },
        {
            **payloads[0],
            "selection": {**payloads[0]["selection"], "module": "os"},
        },
        {**payloads[1], "sample_limit": 1},
        {**payloads[2], "required_tag_ids": []},
        {**payloads[3], "card_version_id": "wrong-alias"},
    ):
        with pytest.raises(ValidationError):
            adapter.validate_python(invalid)


def test_all_http_response_contract_fields_serialize_lower_camel() -> None:
    from course_helper.api import (
        EvidenceResponse,
        HealthResponse,
        JobResponse,
        KnowledgeSummary,
        SessionExchangeResponse,
    )

    now = datetime(2026, 7, 16, tzinfo=timezone.utc)
    evidence = EvidenceResponse.model_validate(
        {
            "evidence_id": "response-evidence",
            "kind": "execution",
            "subject_version_id": None,
            "status": "verified",
            "input_summary": {},
            "output_summary": {},
            "producer": "response-tests",
            "producer_version": "1",
            "started_at": now,
            "finished_at": now,
            "duration_ms": 0,
            "artifacts": [
                {
                    "artifact_id": "artifact-1",
                    "locator": {"root_id": "app-data", "relative_path": "safe.json"},
                    "media_type": "application/json",
                    "content_digest": "a" * 64,
                    "byte_size": 10,
                }
            ],
        }
    )
    response_payloads = (
        SessionExchangeResponse(session_token="t" * 43).model_dump(
            mode="json", by_alias=True
        ),
        HealthResponse(
            service_version="0.1.0",
            schema_version=1,
            database_ready=True,
        ).model_dump(mode="json", by_alias=True),
        KnowledgeSummary(
            source_count=1,
            published_card_count=2,
            review_task_count=3,
            retrieval_mode="fts-degraded",
            index_snapshot_id=None,
            index_snapshot_digest=None,
            index_state="unavailable",
            tag_labels=("RFM",),
            tag_options=({"id": "topic:rfm", "label": "RFM", "dimension": "topic"},),
            updated_at=now,
        ).model_dump(mode="json", by_alias=True),
        JobResponse(result={"hitCount": 0}, evidence=evidence).model_dump(
            mode="json", by_alias=True
        ),
    )

    assert response_payloads[0] == {"sessionToken": "t" * 43}
    assert set(response_payloads[1]) == {
        "serviceVersion",
        "schemaVersion",
        "databaseReady",
    }
    assert set(response_payloads[2]) == {
        "schemaVersion",
        "sourceCount",
        "publishedCardCount",
        "reviewTaskCount",
        "retrievalMode",
        "indexSnapshotId",
        "indexSnapshotDigest",
        "indexState",
        "tagLabels",
        "tagOptions",
        "updatedAt",
    }
    evidence_payload = response_payloads[3]["evidence"]
    assert "evidenceId" in evidence_payload and "evidence_id" not in evidence_payload
    assert evidence_payload["artifacts"][0]["artifactId"] == "artifact-1"
    assert evidence_payload["artifacts"][0]["locator"] == {
        "rootId": "app-data",
        "relativePath": "safe.json",
    }


def test_launch_nonce_is_exact_origin_single_use_and_expires(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from course_helper.session import LaunchSession, SessionRejected

    issued = iter(("n" * 43, "t" * 43, "x" * 43, "y" * 43))
    now = [100.0]
    monkeypatch.setattr("course_helper.session.secrets.token_urlsafe", lambda _: next(issued))
    session = LaunchSession.create(
        allowed_origin="http://127.0.0.1:4173",
        monotonic=lambda: now[0],
    )

    assert len(session.launch_nonce) >= 43
    assert "n" * 43 not in repr(session)
    assert "t" * 43 not in repr(session)
    with pytest.raises(SessionRejected):
        session.exchange("n" * 43, origin="http://localhost:4173")
    assert session.exchange(
        "n" * 43,
        origin="http://127.0.0.1:4173",
    ) == "t" * 43
    with pytest.raises(SessionRejected):
        session.exchange("n" * 43, origin="http://127.0.0.1:4173")

    expired = LaunchSession.create(
        allowed_origin="http://127.0.0.1:4173",
        monotonic=lambda: now[0],
    )
    now[0] = 160.0
    with pytest.raises(SessionRejected):
        expired.exchange("x" * 43, origin="http://127.0.0.1:4173")


def test_launch_nonce_exchange_is_atomic_across_threads() -> None:
    from course_helper.session import LaunchSession, SessionRejected

    origin = "http://127.0.0.1:4173"

    for _round in range(20):
        exchange_barrier = threading.Barrier(20)
        main_thread_id = threading.get_ident()

        def racing_clock() -> float:
            if threading.get_ident() != main_thread_id:
                try:
                    exchange_barrier.wait(timeout=0.05)
                except threading.BrokenBarrierError:
                    pass
            return 100.0

        session = LaunchSession(
            allowed_origin=origin,
            launch_nonce="n" * 43,
            session_token="t" * 43,
            expires_at=160.0,
            monotonic=racing_clock,
        )
        outcomes: list[tuple[str, str]] = []

        def exchange() -> None:
            try:
                outcomes.append(("accepted", session.exchange("n" * 43, origin=origin)))
            except SessionRejected as error:
                outcomes.append(("rejected", str(error)))

        threads = [threading.Thread(target=exchange) for _ in range(20)]
        for thread in threads:
            thread.start()
        for thread in threads:
            thread.join(timeout=2.0)

        assert all(not thread.is_alive() for thread in threads)
        accepted = [value for status, value in outcomes if status == "accepted"]
        rejected = [value for status, value in outcomes if status == "rejected"]
        assert accepted == ["t" * 43]
        assert rejected == ["launch session rejected"] * 19


def _api_client(tmp_path: Path) -> tuple[TestClient, str]:
    from course_helper.api import HelperRuntime, create_app
    from course_helper.jobs import JobOutcome, KnowledgeRetrieveJob, WorkerRuntimeConfig
    from course_helper.session import LaunchSession

    database_path = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database_path):
        pass
    source_root = tmp_path / "sources"
    source_root.mkdir()
    session = LaunchSession.create(allowed_origin="http://127.0.0.1:4173")

    class InlineRunner:
        async def run(
            self,
            job: object,
            *,
            disconnected: object = None,
            session_id: str | None = None,
        ) -> JobOutcome:
            assert isinstance(job, KnowledgeRetrieveJob)
            assert session_id is not None and session_id.startswith("session-")
            now = datetime(2026, 7, 16, tzinfo=timezone.utc)
            return JobOutcome(
                status_code=200,
                result={"hitCount": 0},
                evidence=EvidenceObject(
                    evidence_id="retrieval-api-fixture",
                    kind="retrieval",
                    status="verified",
                    producer="api-tests",
                    started_at=now,
                    finished_at=now,
                    duration_ms=0,
                    checks=(
                        EvidenceCheck(
                            code="bounded-retrieval",
                            status="passed",
                            message="Retrieval completed inside configured ceilings",
                        ),
                    ),
                ),
            )

    runtime = HelperRuntime(
        config=WorkerRuntimeConfig(
            database_path=str(database_path),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(("fixture", str(source_root)),),
        ),
        launch_session=session,
        job_runner=InlineRunner(),
    )
    client = TestClient(create_app(runtime))
    return client, session.launch_nonce


def _authenticated_headers(client: TestClient, nonce: str) -> dict[str, str]:
    origin = "http://127.0.0.1:4173"
    response = client.post(
        "/v1/session/exchange",
        headers={"Origin": origin},
        json={"nonce": nonce},
    )
    assert response.status_code == 200
    token = response.json()["sessionToken"]
    assert len(token) >= 43
    return {"Origin": origin, "X-Course-Session": token}


def test_authenticated_course_projection_reopens_identical_published_bytes(
    tmp_path: Path,
) -> None:
    from datetime import timedelta

    from course_helper.api import HelperRuntime, create_app
    from course_helper.jobs import WorkerRuntimeConfig
    from course_helper.session import LaunchSession
    from course_helper.slide_builder import publish_course_version
    from test_course_publication import _prepare_publication, _request
    from test_network_visuals import NOW as NETWORK_NOW

    value = _prepare_publication(tmp_path)
    outcome = publish_course_version(
        value.catalog,
        _request(value, "api-course-projection", value.visual_placement_ids),
        confirmed_course_version_id=value.confirmed_course_id,
        expected_course_digest=value.confirmed_course_digest,
        visual_placement_ids=value.visual_placement_ids,
        clock=lambda: NETWORK_NOW + timedelta(hours=1),
    )
    value.catalog.close()
    course_id = str(outcome.result_refs["courseVersionId"])
    deck_id = str(outcome.result_refs["slideDeckId"])
    manifest_id = str(outcome.result_refs["runtimeManifestId"])

    class NeverRunner:
        async def run(self, *_args: object, **_kwargs: object) -> object:
            raise AssertionError("projection reopen must not dispatch a worker job")

    session = LaunchSession.create(allowed_origin="http://127.0.0.1:4173")
    client = TestClient(create_app(HelperRuntime(
        config=WorkerRuntimeConfig(
            database_path=str(tmp_path / "publication.sqlite3"),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        ),
        launch_session=session,
        job_runner=NeverRunner(),  # type: ignore[arg-type]
    )))
    path = (
        f"/v1/courses/{course_id}/projection"
        f"?slideDeckId={deck_id}&runtimeManifestId={manifest_id}"
    )
    assert client.get(path).status_code == 401
    headers = _authenticated_headers(client, session.launch_nonce)
    first = client.get(path, headers=headers)
    second = client.get(path, headers=headers)
    assert first.status_code == second.status_code == 200
    assert first.content == second.content
    payload = first.json()
    assert payload["courseVersionId"] == course_id
    assert payload["slideDeck"]["versionId"] == deck_id
    assert payload["runtimeManifest"]["versionId"] == manifest_id
    assert payload["runtimeManifest"]["slideDeckDigest"] == payload["slideDeck"]["contentDigest"]
    assert payload["requirement"]["requirementId"] == payload["outline"]["requirementId"]
    assert client.get(
        f"/v1/courses/{course_id}/projection?slideDeckId=wrong&runtimeManifestId={manifest_id}",
        headers=headers,
    ).status_code == 404


def test_api_requires_exact_origin_and_session_and_nonce_is_single_use(
    tmp_path: Path,
) -> None:
    client, nonce = _api_client(tmp_path)

    assert client.post(
        "/v1/jobs",
        json={"type": "knowledge_retrieve", "query": "RFM"},
    ).status_code == 401
    for origin_headers in (
        {"Origin": "http://localhost:4173"},
        {"Origin": "http://127.0.0.1:4173/"},
        {},
    ):
        rejected = client.post(
            "/v1/session/exchange",
            headers=origin_headers,
            json={"nonce": nonce},
        )
        assert rejected.status_code == 401
        assert nonce not in rejected.text
    wrong_nonce = client.post(
        "/v1/session/exchange",
        headers={"Origin": "http://127.0.0.1:4173"},
        json={"nonce": "wrong-launch-nonce"},
    )
    invalid_body = client.post(
        "/v1/session/exchange",
        headers={"Origin": "http://127.0.0.1:4173"},
        json={"nonce": nonce, "sessionToken": "raw-token-must-not-echo"},
    )
    assert wrong_nonce.status_code == 401
    assert invalid_body.status_code == 422
    assert "wrong-launch-nonce" not in wrong_nonce.text
    assert nonce not in invalid_body.text
    assert "raw-token-must-not-echo" not in invalid_body.text

    headers = _authenticated_headers(client, nonce)
    assert client.post(
        "/v1/session/exchange",
        headers={"Origin": headers["Origin"]},
        json={"nonce": nonce},
    ).status_code == 401
    assert client.post(
        "/v1/jobs",
        headers={**headers, "Origin": "http://localhost:4173"},
        json={"type": "knowledge_retrieve", "query": "RFM"},
    ).status_code == 401


def test_cors_allows_only_the_exact_configured_origin(tmp_path: Path) -> None:
    client, _nonce = _api_client(tmp_path)

    allowed = client.options(
        "/v1/jobs",
        headers={
            "Origin": "http://127.0.0.1:4173",
            "Access-Control-Request-Method": "POST",
            "Access-Control-Request-Headers": "X-Course-Session, Content-Type",
        },
    )
    disallowed = client.options(
        "/v1/jobs",
        headers={
            "Origin": "http://localhost:4173",
            "Access-Control-Request-Method": "POST",
        },
    )

    assert allowed.headers["access-control-allow-origin"] == "http://127.0.0.1:4173"
    assert "access-control-allow-origin" not in disallowed.headers


def test_cors_allows_only_the_explicit_upload_header(tmp_path: Path) -> None:
    client, _nonce = _api_client(tmp_path)

    allowed = client.options(
        "/v1/uploads",
        headers={
            "Origin": "http://127.0.0.1:4173",
            "Access-Control-Request-Method": "POST",
            "Access-Control-Request-Headers": (
                "X-Course-Session, X-Upload-Name, Content-Type"
            ),
        },
    )
    rejected = client.options(
        "/v1/uploads",
        headers={
            "Origin": "http://127.0.0.1:4173",
            "Access-Control-Request-Method": "POST",
            "Access-Control-Request-Headers": "X-Course-Session, X-Source-Path",
        },
    )

    assert allowed.status_code == 200
    assert allowed.headers["access-control-allow-origin"] == "http://127.0.0.1:4173"
    assert rejected.status_code == 400
    assert str(tmp_path) not in rejected.text


def test_authenticated_streaming_upload_returns_only_safe_opaque_metadata(
    tmp_path: Path,
) -> None:
    client, nonce = _api_client(tmp_path)
    content = b"# Governed upload\n\nOne source.\n"
    unauthenticated = client.post(
        "/v1/uploads",
        headers={"X-Upload-Name": "lesson.md", "Content-Type": "text/markdown"},
        content=content,
    )
    headers = _authenticated_headers(client, nonce)
    response = client.post(
        "/v1/uploads",
        headers={
            **headers,
            "X-Upload-Name": "lesson.md",
            "Content-Type": "text/markdown; charset=utf-8",
        },
        content=content,
    )

    assert unauthenticated.status_code == 401
    assert response.status_code == 201
    assert response.headers["cache-control"] == "no-store"
    payload = response.json()
    assert set(payload) == {
        "schemaVersion",
        "uploadId",
        "safeName",
        "sourceKind",
        "mediaType",
        "byteSize",
        "contentDigest",
        "state",
        "expiresAt",
    }
    assert payload["uploadId"].startswith("upload-")
    assert payload["safeName"] == "lesson.md"
    assert payload["byteSize"] == len(content)
    assert payload["contentDigest"] == hashlib.sha256(content).hexdigest()
    assert payload["state"] == "available"
    serialized = response.text + unauthenticated.text
    assert str(tmp_path) not in serialized
    assert headers["X-Course-Session"] not in serialized
    assert content.decode() not in serialized

    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database) as catalog:
        row = catalog.connection.execute(
            "SELECT safe_name, byte_size, content_digest, session_digest "
            "FROM governed_uploads "
            "WHERE upload_id = ?",
            (payload["uploadId"],),
        ).fetchone()
    owner_id = "session-" + hashlib.sha256(
        headers["X-Course-Session"].encode()
    ).hexdigest()
    assert tuple(row) == (
        "lesson.md",
        len(content),
        hashlib.sha256(content).hexdigest(),
        hashlib.sha256(owner_id.encode()).hexdigest(),
    )
    assert (tmp_path / "app-data" / "uploads" / f"{payload['uploadId']}.blob").read_bytes() == content


@pytest.mark.parametrize(
    ("headers", "content", "status"),
    (
        ({"Content-Type": "text/markdown"}, b"# Missing name\n", 422),
        (
            {"X-Upload-Name": "lesson.md", "Content-Type": "application/octet-stream"},
            b"# Wrong media\n",
            422,
        ),
        (
            {
                "X-Upload-Name": "lesson.md",
                "Content-Type": "text/markdown",
                "Content-Length": str(20 * 1024 * 1024 + 1),
            },
            b"x",
            413,
        ),
        (
            {
                "X-Upload-Name": "lesson.md",
                "Content-Type": "text/markdown",
                "Content-Length": "99",
            },
            b"short",
            422,
        ),
    ),
)
def test_upload_boundary_rejects_invalid_headers_length_and_media_without_residue(
    tmp_path: Path,
    headers: dict[str, str],
    content: bytes,
    status: int,
) -> None:
    client, nonce = _api_client(tmp_path)
    authenticated = _authenticated_headers(client, nonce)
    response = client.post(
        "/v1/uploads",
        headers={**authenticated, **headers},
        content=content,
    )

    assert response.status_code == status
    assert set(response.json()) == {"error"}
    assert set(response.json()["error"]) == {"code", "message"}
    serialized = response.text
    assert str(tmp_path) not in serialized
    assert authenticated["X-Course-Session"] not in serialized
    assert not tuple((tmp_path / "app-data").rglob("*.tmp"))
    with KnowledgeCatalog.open(tmp_path / "knowledge.db") as catalog:
        assert catalog.connection.execute(
            "SELECT count(*) FROM governed_uploads"
        ).fetchone()[0] == 0


def _seed_governed_inventory_source(
    tmp_path: Path,
    *,
    suffix: str,
    session_id: str,
    created_at: datetime,
) -> str:
    from course_helper.domain.common import ActorRef
    from course_helper.operations import OperationRequest
    from course_helper.uploads import (
        UploadStore,
        import_promotion_request_digest,
        import_start_request_digest,
    )

    content = f"# Inventory source {suffix}\n".encode()
    owner_id = "session-" + hashlib.sha256(session_id.encode("utf-8")).hexdigest()
    with KnowledgeCatalog.open(tmp_path / "knowledge.db") as catalog:
        store = UploadStore(catalog, tmp_path / "app-data")
        upload = store.create_upload(
            (content,),
            file_name=f"source-{suffix}.md",
            media_type="text/markdown",
            byte_size_hint=len(content),
            session_id=owner_id,
            clock=lambda: created_at,
        )
        actor = ActorRef(actor_type="human", actor_id="api-inventory-author")
        start_digest = import_start_request_digest(upload.upload_id, upload.content_digest)
        started = store.start_import(
            OperationRequest(
                operation_id=f"api-inventory-start-{suffix}",
                request_digest=start_digest,
                actor=actor,
                session_id=owner_id,
            ),
            upload_id=upload.upload_id,
            expected_content_digest=upload.content_digest,
            clock=lambda: created_at,
        )
        import_id = str(started.result_refs["importId"])
        promotion_digest = import_promotion_request_digest(import_id, upload.content_digest)
        promoted = store.promote_import(
            OperationRequest(
                operation_id=f"api-inventory-promote-{suffix}",
                request_digest=promotion_digest,
                actor=actor,
                session_id=owner_id,
            ),
            import_id=import_id,
            expected_content_digest=upload.content_digest,
            clock=lambda: created_at,
        )
    return str(promoted.result_refs["sourceVersionId"])


def test_authenticated_source_inventory_is_bounded_paginated_and_path_free(
    tmp_path: Path,
) -> None:
    client, nonce = _api_client(tmp_path)
    headers = _authenticated_headers(client, nonce)
    session_id = headers["X-Course-Session"]
    now = datetime(2026, 7, 19, 4, 0, tzinfo=timezone.utc)
    expected_ids = tuple(
        _seed_governed_inventory_source(
            tmp_path,
            suffix=str(index),
            session_id=session_id,
            created_at=now.replace(second=index),
        )
        for index in range(3)
    )

    unauthenticated = client.get("/v1/knowledge/sources")
    first = client.get("/v1/knowledge/sources", headers=headers, params={"limit": "2"})
    second = client.get(
        "/v1/knowledge/sources",
        headers=headers,
        params={"limit": "2", "cursor": first.json()["nextCursor"]},
    )
    injected = client.get(
        "/v1/knowledge/sources",
        headers=headers,
        params={"limit": "2", "sourcePath": str(tmp_path / "private")},
    )

    assert unauthenticated.status_code == 401
    assert first.status_code == second.status_code == 200
    assert first.headers["cache-control"] == second.headers["cache-control"] == "no-store"
    assert injected.status_code == 422
    assert [item["sourceVersionId"] for item in first.json()["items"]] == list(
        expected_ids[:2]
    )
    assert [item["sourceVersionId"] for item in second.json()["items"]] == list(
        expected_ids[2:]
    )
    assert first.json()["nextCursor"].startswith("inventory-cursor-")
    assert second.json()["nextCursor"] is None
    assert set(first.json()["items"][0]) == {
        "schemaVersion",
        "sourceId",
        "sourceVersionId",
        "displayName",
        "sourceKind",
        "mediaType",
        "byteSize",
        "contentDigest",
        "status",
    }
    serialized = first.text + second.text + injected.text
    assert str(tmp_path) not in serialized
    assert session_id not in serialized
    assert "governed-upload" not in serialized
    assert "relativePath" not in serialized
    assert "# Inventory" not in serialized


@pytest.mark.parametrize(
    "params",
    (
        {"limit": "0"},
        {"limit": "101"},
        {"limit": "01"},
        {"limit": "two"},
        {"cursor": "../private"},
        (("limit", "2"), ("limit", "3")),
    ),
)
def test_source_inventory_rejects_noncanonical_or_repeated_query_parameters(
    tmp_path: Path,
    params: object,
) -> None:
    client, nonce = _api_client(tmp_path)
    response = client.get(
        "/v1/knowledge/sources",
        headers=_authenticated_headers(client, nonce),
        params=params,
    )

    assert response.status_code == 422
    assert str(tmp_path) not in response.text


def test_api_rejects_unknown_shell_and_extra_fields_before_dispatch(tmp_path: Path) -> None:
    client, nonce = _api_client(tmp_path)
    headers = _authenticated_headers(client, nonce)

    raw_command = "whoami --secret-command"
    shell = client.post(
        "/v1/jobs",
        headers=headers,
        json={"type": "shell", "command": raw_command},
    )
    extra = client.post(
        "/v1/jobs",
        headers=headers,
        json={"type": "knowledge_retrieve", "query": "RFM", "command": raw_command},
    )
    publish_json = client.post(
        "/v1/jobs",
        headers=headers,
        json={"type": "knowledge_publish", "card": {"title": "unsafe"}},
    )

    assert shell.status_code == 422
    assert extra.status_code == 422
    assert publish_json.status_code == 422
    serialized = shell.text + extra.text + publish_json.text
    assert raw_command not in serialized
    assert "unsafe" not in serialized


def test_typed_retrieval_returns_lower_camel_structured_evidence(tmp_path: Path) -> None:
    client, nonce = _api_client(tmp_path)
    response = client.post(
        "/v1/jobs",
        headers=_authenticated_headers(client, nonce),
        json={"type": "knowledge_retrieve", "query": "RFM", "limit": 5},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["evidence"]["kind"] == "retrieval"
    assert payload["evidence"]["evidenceId"] == "retrieval-api-fixture"
    assert payload["result"] == {"hitCount": 0}
    assert response.headers["cache-control"] == "no-store"


def test_api_injects_derived_session_owner_into_import_job_without_echoing_token(
    tmp_path: Path,
) -> None:
    from course_helper.api import HelperRuntime, create_app
    from course_helper.jobs import (
        JobOutcome,
        KnowledgeImportStartJob,
        WorkerRuntimeConfig,
    )
    from course_helper.session import LaunchSession
    from course_helper.uploads import import_start_request_digest

    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    session = LaunchSession.create(allowed_origin="http://127.0.0.1:4173")
    captured: dict[str, object] = {}

    class CaptureRunner:
        async def run(
            self,
            job: object,
            *,
            disconnected: object,
            session_id: str | None = None,
        ) -> JobOutcome:
            assert isinstance(job, KnowledgeImportStartJob)
            captured.update(job=job, session_id=session_id)
            now = datetime(2026, 7, 19, 5, 0, tzinfo=timezone.utc)
            return JobOutcome(
                status_code=200,
                result={"operationStatus": "committed"},
                evidence=EvidenceObject(
                    evidence_id="import-api-capture",
                    kind="execution",
                    status="verified",
                    producer="api-tests",
                    started_at=now,
                    finished_at=now,
                ),
            )

    client = TestClient(
        create_app(
            HelperRuntime(
                config=WorkerRuntimeConfig(
                    database_path=str(database),
                    app_data_path=str(tmp_path / "app-data"),
                    source_roots=(),
                ),
                launch_session=session,
                job_runner=CaptureRunner(),
            )
        )
    )
    headers = _authenticated_headers(client, session.launch_nonce)
    upload_id = "upload-" + "1" * 32
    content_digest = "2" * 64
    response = client.post(
        "/v1/jobs",
        headers=headers,
        json={
            "type": "knowledge_import_start",
            "uploadId": upload_id,
            "expectedContentDigest": content_digest,
            "operationId": "api-import-operation",
            "requestDigest": import_start_request_digest(upload_id, content_digest),
            "actor": {"actorType": "human", "actorId": "api-author"},
        },
    )

    expected_owner = "session-" + hashlib.sha256(
        headers["X-Course-Session"].encode()
    ).hexdigest()
    assert response.status_code == 200
    assert captured["session_id"] == expected_owner
    serialized = response.text + repr(captured["job"])
    assert headers["X-Course-Session"] not in serialized
    assert str(tmp_path) not in serialized


def test_health_and_authenticated_summary_are_path_and_secret_free(tmp_path: Path) -> None:
    client, nonce = _api_client(tmp_path)
    health = client.get("/health")
    headers = _authenticated_headers(client, nonce)
    summary = client.get(
        "/v1/knowledge/summary",
        headers=headers,
    )
    path_injection = client.get(
        "/v1/knowledge/summary",
        headers=headers,
        params={"sourcePath": str(tmp_path / "secret")},
    )

    assert health.status_code == 200
    assert health.json() == {
        "serviceVersion": "0.1.0",
        "schemaVersion": 7,
        "databaseReady": True,
    }
    assert summary.status_code == 200
    assert summary.json()["schemaVersion"] == 1
    assert summary.json()["sourceCount"] == 0
    assert summary.json()["retrievalMode"] == "fts-degraded"
    assert path_injection.status_code == 422
    serialized = health.text + summary.text + path_injection.text
    assert str(tmp_path) not in serialized
    assert nonce not in serialized


def test_authenticated_summary_counts_only_unsuspended_projected_publications(
    tmp_path: Path,
) -> None:
    from course_helper.cards import publish_card
    from course_helper.domain.knowledge import KnowledgeCardVersion
    from course_helper.lifecycle import append_card_lifecycle_event

    database = tmp_path / "knowledge.db"
    card_version_id = _seed_publish_candidate(database, source_backed=True, name="summary")
    with KnowledgeCatalog.open(database) as catalog:
        payload = catalog.connection.execute(
            "SELECT payload_json FROM cards WHERE version_id = ?",
            (card_version_id,),
        ).fetchone()[0]
        published = publish_card(KnowledgeCardVersion.model_validate_json(payload), catalog)
        append_card_lifecycle_event(
            catalog.connection,
            card_version_id=published.version_id,
            event_id="suspend-summary-publication",
            request_digest="3" * 64,
            event_type="suspend",
            occurred_at=datetime(2026, 7, 17, tzinfo=timezone.utc),
            actor_id="api-tests",
        )

    client, nonce = _api_client(tmp_path)
    response = client.get(
        "/v1/knowledge/summary",
        headers=_authenticated_headers(client, nonce),
    )

    assert response.status_code == 200
    assert response.json()["publishedCardCount"] == 0


class _FakeQueue:
    def get(self, timeout: float | None = None) -> object:
        raise AssertionError("timeout and cancellation must not read an empty child queue")

    def close(self) -> None:
        pass

    def join_thread(self) -> None:
        pass


class _RecordingQueue:
    def __init__(self) -> None:
        self.items: list[dict[str, object]] = []

    def put(self, value: dict[str, object]) -> None:
        self.items.append(value)


class _AliveProcess:
    def __init__(self) -> None:
        self.started = False
        self.terminated = False
        self.joined = False
        self.exitcode: int | None = None

    def start(self) -> None:
        self.started = True

    def is_alive(self) -> bool:
        return not self.terminated

    def terminate(self) -> None:
        self.terminated = True
        self.exitcode = -15

    def join(self, timeout: float | None = None) -> None:
        self.joined = True


class _FakeContext:
    def __init__(self) -> None:
        self.process = _AliveProcess()
        self.process_calls = 0
        self.target: object = None
        self.args: tuple[object, ...] = ()

    def Queue(self) -> _FakeQueue:
        return _FakeQueue()

    def Process(self, *, target: object, args: tuple[object, ...]) -> _AliveProcess:
        self.process_calls += 1
        self.target = target
        self.args = args
        return self.process


class _ClosableQueue(_FakeQueue):
    def __init__(self) -> None:
        self.closed = False
        self.joined = False

    def close(self) -> None:
        self.closed = True

    def join_thread(self) -> None:
        self.joined = True


class _QueueConstructorFailureContext:
    def __init__(self, message: str) -> None:
        self.message = message
        self.process_calls = 0

    def Queue(self) -> _FakeQueue:
        raise OSError(self.message)

    def Process(self, *, target: object, args: tuple[object, ...]) -> _AliveProcess:
        self.process_calls += 1
        raise AssertionError("process construction must not follow queue failure")


class _ProcessConstructorFailureContext:
    def __init__(self, message: str) -> None:
        self.message = message
        self.queue = _ClosableQueue()
        self.process_calls = 0

    def Queue(self) -> _ClosableQueue:
        return self.queue

    def Process(self, *, target: object, args: tuple[object, ...]) -> _AliveProcess:
        self.process_calls += 1
        raise OSError(self.message)


class _StartFailureProcess:
    def __init__(self, message: str) -> None:
        self.message = message
        self.start_calls = 0
        self.terminate_calls = 0
        self.join_calls = 0
        self.close_calls = 0
        self.exitcode: int | None = None

    def start(self) -> None:
        self.start_calls += 1
        raise OSError(self.message)

    def is_alive(self) -> bool:
        raise AssertionError("an unstarted process must not be inspected")

    def terminate(self) -> None:
        self.terminate_calls += 1

    def join(self, timeout: float | None = None) -> None:
        self.join_calls += 1

    def close(self) -> None:
        self.close_calls += 1


class _StartFailureContext:
    def __init__(self, message: str) -> None:
        self.queue = _ClosableQueue()
        self.process = _StartFailureProcess(message)

    def Queue(self) -> _ClosableQueue:
        return self.queue

    def Process(self, *, target: object, args: tuple[object, ...]) -> _StartFailureProcess:
        return self.process


async def _connected() -> bool:
    return False


def _run_retrieval_with_context(
    tmp_path: Path,
    context: object,
) -> tuple[object, tuple[Path, ...]]:
    from course_helper.jobs import BoundedJobRunner, KnowledgeRetrieveJob, WorkerRuntimeConfig

    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        ),
        mp_context=context,
    )
    outcome = asyncio.run(
        runner.run(
            KnowledgeRetrieveJob(type="knowledge_retrieve", query="RFM"),
            disconnected=_connected,
        )
    )
    receipts = tuple((tmp_path / "app-data" / "job-evidence").glob("*.json"))
    return outcome, receipts


def test_queue_constructor_failure_returns_sanitized_evidence_without_spawning(
    tmp_path: Path,
) -> None:
    secret = str(tmp_path / "private-queue-material")
    context = _QueueConstructorFailureContext(secret)

    outcome, receipts = _run_retrieval_with_context(tmp_path, context)

    assert outcome.status_code == 500
    assert outcome.evidence.checks[0].code == "job-failed"
    assert outcome.evidence.checks[0].details["reason_code"] == "OSError"
    assert context.process_calls == 0
    assert len(receipts) == 1
    serialized = outcome.evidence.model_dump_json() + receipts[0].read_text(encoding="utf-8")
    assert secret not in serialized
    assert str(tmp_path) not in serialized


def test_process_constructor_failure_closes_queue_and_returns_sanitized_evidence(
    tmp_path: Path,
) -> None:
    secret = str(tmp_path / "private-process-material")
    context = _ProcessConstructorFailureContext(secret)

    outcome, receipts = _run_retrieval_with_context(tmp_path, context)

    assert outcome.status_code == 500
    assert outcome.evidence.checks[0].code == "job-failed"
    assert outcome.evidence.checks[0].details["reason_code"] == "OSError"
    assert context.process_calls == 1
    assert context.queue.closed is True
    assert context.queue.joined is True
    assert len(receipts) == 1
    serialized = outcome.evidence.model_dump_json() + receipts[0].read_text(encoding="utf-8")
    assert secret not in serialized
    assert str(tmp_path) not in serialized


def test_process_start_failure_closes_resources_without_terminate_or_join(
    tmp_path: Path,
) -> None:
    secret = str(tmp_path / "private-start-material")
    context = _StartFailureContext(secret)

    outcome, receipts = _run_retrieval_with_context(tmp_path, context)

    assert outcome.status_code == 500
    assert outcome.evidence.checks[0].code == "job-failed"
    assert outcome.evidence.checks[0].details["reason_code"] == "OSError"
    assert context.process.start_calls == 1
    assert context.process.terminate_calls == 0
    assert context.process.join_calls == 0
    assert context.process.close_calls == 1
    assert context.queue.closed is True
    assert context.queue.joined is True
    assert len(receipts) == 1
    serialized = outcome.evidence.model_dump_json() + receipts[0].read_text(encoding="utf-8")
    assert secret not in serialized
    assert str(tmp_path) not in serialized


def test_preflight_size_violation_returns_413_without_spawning(tmp_path: Path) -> None:
    from course_helper.jobs import BoundedJobRunner, SourceIngestJob, WorkerRuntimeConfig

    source_root = tmp_path / "sources"
    source_root.mkdir()
    oversized = source_root / "oversized.pptx"
    with oversized.open("wb") as output:
        output.truncate(512 * 1024 * 1024 + 1)
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    context = _FakeContext()
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(("fixture", str(source_root)),),
    )
    runner = BoundedJobRunner(config, mp_context=context)
    job = SourceIngestJob.model_validate(
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "oversized.pptx"},
            "selection": {"slideNumbers": [1]},
        }
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 413
    assert context.process_calls == 0
    assert outcome.evidence.checks[0].code == "job-preflight"
    assert str(tmp_path) not in outcome.evidence.model_dump_json()


def test_timeout_terminates_and_joins_spawn_child_with_failure_evidence(
    tmp_path: Path,
) -> None:
    from course_helper.jobs import (
        BoundedJobRunner,
        KnowledgeRetrieveJob,
        WorkerRuntimeConfig,
        _spawn_job_entry,
    )

    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    context = _FakeContext()
    ticks = iter((0.0, 6.0))
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        ),
        mp_context=context,
        monotonic=lambda: next(ticks),
        poll_interval=0,
    )
    job = KnowledgeRetrieveJob(
        type="knowledge_retrieve",
        query="RFM",
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 504
    assert outcome.evidence.checks[0].code == "job-timeout"
    assert context.process.started is True
    assert context.process.terminated is True
    assert context.process.joined is True
    assert context.target is _spawn_job_entry
    assert len(context.args) == 3
    assert isinstance(context.args[0], dict)
    assert context.args[1] == runner.config


def test_session_secrets_never_enter_worker_config_args_evidence_repr_or_logs(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    caplog: pytest.LogCaptureFixture,
) -> None:
    from course_helper.jobs import BoundedJobRunner, KnowledgeRetrieveJob, WorkerRuntimeConfig
    from course_helper.session import LaunchSession

    nonce = "nonce-secret-material-that-must-stay-in-session-123456"
    token = "token-secret-material-that-must-stay-in-session-654321"
    issued = iter((nonce, token))
    monkeypatch.setattr("course_helper.session.secrets.token_urlsafe", lambda _: next(issued))
    session = LaunchSession.create(allowed_origin="http://127.0.0.1:4173")
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    context = _FakeContext()
    ticks = iter((0.0, 6.0))
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(),
    )
    runner = BoundedJobRunner(
        config,
        mp_context=context,
        monotonic=lambda: next(ticks),
        poll_interval=0,
    )

    outcome = asyncio.run(
        runner.run(
            KnowledgeRetrieveJob(type="knowledge_retrieve", query="RFM"),
            disconnected=_connected,
        )
    )
    receipts = tuple((tmp_path / "app-data" / "job-evidence").glob("*.json"))
    combined = "\n".join(
        (
            repr(session),
            repr(config),
            repr(context.args),
            outcome.evidence.model_dump_json(),
            *(receipt.read_text(encoding="utf-8") for receipt in receipts),
            caplog.text,
        )
    )

    assert outcome.status_code == 504
    assert nonce not in combined
    assert token not in combined
    assert "launch_nonce" not in repr(config)
    assert "session_token" not in repr(config)


def test_disconnect_terminates_child_and_persists_sanitized_cancel_evidence(
    tmp_path: Path,
) -> None:
    from course_helper.jobs import BoundedJobRunner, KnowledgeRetrieveJob, WorkerRuntimeConfig

    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    context = _FakeContext()
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        ),
        mp_context=context,
        poll_interval=0,
    )

    async def disconnected() -> bool:
        return True

    outcome = asyncio.run(
        runner.run(
            KnowledgeRetrieveJob(type="knowledge_retrieve", query="RFM"),
            disconnected=disconnected,
        )
    )

    assert outcome.status_code == 499
    assert context.process.terminated is True
    assert context.process.joined is True
    receipts = tuple((tmp_path / "app-data" / "job-evidence").glob("*.json"))
    assert len(receipts) == 1
    persisted = json.loads(receipts[0].read_text(encoding="utf-8"))
    assert persisted["checks"][0]["code"] == "job-cancelled"
    assert str(tmp_path) not in json.dumps(persisted)


def test_real_spawn_dispatches_only_the_module_level_allowlisted_jobs(
    tmp_path: Path,
) -> None:
    from course_helper.cards import seed_vocabulary
    from course_helper.jobs import (
        BoundedJobRunner,
        KnowledgeRetrieveJob,
        WorkerRuntimeConfig,
        _ALLOWLISTED_HANDLERS,
    )

    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database) as catalog:
        seed_vocabulary(catalog)
    assert tuple(_ALLOWLISTED_HANDLERS) == (
        "source_ingest",
        "dataset_profile",
        "knowledge_retrieve",
        "knowledge_publish",
        "knowledge_import_start",
        "knowledge_import_status",
        "knowledge_import_cancel",
        "operation_status",
        "knowledge_review_list",
        "knowledge_review_detail",
        "knowledge_upgrade_list",
        "knowledge_review_resolve",
        "knowledge_card_publish",
        "knowledge_upgrade_resolve",
        "knowledge_index",
        "course_compose",
        "course_outline_confirm",
        "chart_build",
        "visual_search",
        "visual_acquire",
        "visual_revalidate",
        "course_visual_attach",
        "course_visual_detach",
        "course_validate",
        "course_publish",
    )
    assert all(handler.__module__ == "course_helper.jobs" for handler in _ALLOWLISTED_HANDLERS.values())
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        )
    )

    outcome = asyncio.run(
        runner.run(
            KnowledgeRetrieveJob(type="knowledge_retrieve", query="RFM", limit=5),
            disconnected=_connected,
        )
    )

    assert outcome.status_code == 200
    assert outcome.evidence.kind == "retrieval"
    assert tuple(check.code for check in outcome.evidence.checks[:3]) == (
        "job-ceilings",
        "job-exit",
        "job-verification",
    )
    assert outcome.evidence.output_summary["exit_code"] == 0
    assert outcome.result["schemaVersion"] == 1
    assert str(tmp_path) not in outcome.evidence.model_dump_json()


def test_parquet_size_preflight_rejects_before_spawn(tmp_path: Path) -> None:
    from course_helper.jobs import BoundedJobRunner, DatasetProfileJob, WorkerRuntimeConfig

    source_root = tmp_path / "sources"
    source_root.mkdir()
    oversized = source_root / "oversized.parquet"
    with oversized.open("wb") as output:
        output.truncate(1024 * 1024 * 1024 + 1)
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    context = _FakeContext()
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(("fixture", str(source_root)),),
        ),
        mp_context=context,
    )
    job = DatasetProfileJob.model_validate(
        {
            "type": "dataset_profile",
            "locator": {"rootId": "fixture", "relativePath": "oversized.parquet"},
            "sampleLimit": 20,
        }
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 413
    assert context.process_calls == 0
    assert outcome.evidence.checks[0].details["reason_code"] == "dataset-too-large"


def test_real_spawn_profiles_synthetic_parquet_through_shared_profiler(
    tmp_path: Path,
) -> None:
    from course_helper.jobs import BoundedJobRunner, DatasetProfileJob, WorkerRuntimeConfig

    source_root = tmp_path / "sources"
    source_root.mkdir()
    parquet_path = source_root / "customers.parquet"
    with duckdb.connect(database=":memory:") as connection:
        connection.execute(
            "CREATE TABLE customers AS SELECT value AS record_id, "
            "'private road ' || value::VARCHAR AS homeAddress "
            "FROM range(1, 4) AS rows(value)"
        )
        connection.table("customers").write_parquet(str(parquet_path))
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(("fixture", str(source_root)),),
        )
    )
    job = DatasetProfileJob.model_validate(
        {
            "type": "dataset_profile",
            "locator": {"rootId": "fixture", "relativePath": "customers.parquet"},
            "sampleLimit": 20,
        }
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 200
    assert outcome.evidence.kind == "dataset-profile"
    assert outcome.result["rowCount"] == 3
    assert outcome.result["sampleCount"] == 3
    assert "private road" not in outcome.evidence.model_dump_json()
    with KnowledgeCatalog.open(database) as catalog:
        stored = catalog.connection.execute("SELECT payload_json FROM datasets").fetchone()
    assert stored is not None
    assert "private road" not in stored[0]


def test_real_spawn_ingests_synthetic_markdown_and_persists_traceable_rows(
    tmp_path: Path,
) -> None:
    from course_helper.jobs import BoundedJobRunner, SourceIngestJob, WorkerRuntimeConfig

    source_root = tmp_path / "sources"
    source_root.mkdir()
    (source_root / "lesson.md").write_text(
        "# RFM lesson\n\nRecency, frequency, and monetary evidence.\n",
        encoding="utf-8",
    )
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(("fixture", str(source_root)),),
        )
    )
    job = SourceIngestJob.model_validate(
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "lesson.md"},
            "selection": {"headingSelectors": ["RFM lesson"]},
        }
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 200
    assert outcome.evidence.kind == "extraction"
    assert outcome.result["chunkCount"] == 1
    with KnowledgeCatalog.open(database) as catalog:
        counts = tuple(
            catalog.connection.execute(f"SELECT count(*) FROM {table}").fetchone()[0]
            for table in ("sources", "chunks", "evidence")
        )
    assert counts == (1, 1, 1)


def test_real_spawn_ingest_failure_is_sanitized_and_leaves_no_catalog_rows(
    tmp_path: Path,
) -> None:
    from course_helper.jobs import BoundedJobRunner, SourceIngestJob, WorkerRuntimeConfig

    source_root = tmp_path / "sources"
    source_root.mkdir()
    (source_root / "broken.md").write_bytes(b"\xff\xfe\xfa")
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(("fixture", str(source_root)),),
        )
    )
    job = SourceIngestJob.model_validate(
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "broken.md"},
            "selection": {"headingSelectors": []},
        }
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 500
    assert outcome.evidence.checks[0].code == "job-failed"
    assert str(tmp_path) not in outcome.evidence.model_dump_json()
    with KnowledgeCatalog.open(database) as catalog:
        assert catalog.connection.execute("SELECT count(*) FROM sources").fetchone()[0] == 0
        assert catalog.connection.execute("SELECT count(*) FROM chunks").fetchone()[0] == 0


def _catalog_bundle_snapshot(database: Path) -> dict[str, tuple[tuple[object, ...], ...]]:
    with KnowledgeCatalog.open(database) as catalog:
        return {
            table: tuple(catalog.connection.execute(f"SELECT * FROM {table} ORDER BY rowid"))
            for table in ("sources", "chunks", "visuals", "datasets", "evidence")
        }


def _write_pptx_ingest_bundle(source_root: Path) -> Path:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    image_path = source_root / "fixture.png"
    Image.new("RGB", (2, 2), color=(36, 99, 235)).save(image_path)
    presentation = Presentation()
    for slide_number in (1, 2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
        textbox.text_frame.text = f"Synthetic slide {slide_number}"
        if slide_number == 1:
            slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.75), Inches(1), Inches(1))
    path = source_root / "bundle.pptx"
    presentation.save(path)
    return path


@pytest.mark.parametrize(
    ("trigger_name", "trigger_sql"),
    (
        (
            "reject_second_chunk",
            """
            CREATE TRIGGER reject_second_chunk
            BEFORE INSERT ON chunks
            WHEN NEW.ordinal = 1
            BEGIN
                SELECT RAISE(ABORT, 'forced second chunk failure');
            END
            """,
        ),
        (
            "reject_visual",
            """
            CREATE TRIGGER reject_visual
            BEFORE INSERT ON visuals
            BEGIN
                SELECT RAISE(ABORT, 'forced visual failure');
            END
            """,
        ),
    ),
)
def test_source_ingest_rolls_back_the_entire_bundle_after_late_trigger_failure(
    tmp_path: Path,
    trigger_name: str,
    trigger_sql: str,
) -> None:
    from course_helper.jobs import SourceIngestJob, WorkerRuntimeConfig, _run_source_ingest

    source_root = tmp_path / "sources"
    source_root.mkdir()
    _write_pptx_ingest_bundle(source_root)
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database) as catalog:
        catalog.connection.execute(trigger_sql)
        catalog.connection.commit()
    before = _catalog_bundle_snapshot(database)
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(("fixture", str(source_root)),),
    )
    job = SourceIngestJob.model_validate(
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "bundle.pptx"},
            "selection": {"slideNumbers": []},
        }
    )

    with pytest.raises(sqlite3.IntegrityError, match=f"forced {trigger_name.removeprefix('reject_').replace('_', ' ')} failure"):
        _run_source_ingest(job, config)

    assert _catalog_bundle_snapshot(database) == before
    with KnowledgeCatalog.open(database) as catalog:
        assert not catalog.connection.in_transaction


def test_dataset_profile_rolls_back_dataset_when_evidence_insert_fails(tmp_path: Path) -> None:
    from course_helper.jobs import DatasetProfileJob, WorkerRuntimeConfig, _run_dataset_profile

    source_root = tmp_path / "sources"
    source_root.mkdir()
    (source_root / "metrics.csv").write_text("metric,value\nquality,1\n", encoding="utf-8")
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database) as catalog:
        catalog.connection.execute(
            """
            CREATE TRIGGER reject_dataset_evidence
            BEFORE INSERT ON evidence
            WHEN NEW.kind = 'dataset-profile'
            BEGIN
                SELECT RAISE(ABORT, 'forced dataset evidence failure');
            END
            """
        )
        catalog.connection.commit()
    before = _catalog_bundle_snapshot(database)
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(("fixture", str(source_root)),),
    )
    job = DatasetProfileJob.model_validate(
        {
            "type": "dataset_profile",
            "locator": {"rootId": "fixture", "relativePath": "metrics.csv"},
            "sampleLimit": 20,
        }
    )

    with pytest.raises(sqlite3.IntegrityError, match="forced dataset evidence failure"):
        _run_dataset_profile(job, config)

    assert _catalog_bundle_snapshot(database) == before
    with KnowledgeCatalog.open(database) as catalog:
        assert not catalog.connection.in_transaction


def _seed_publish_candidate(
    database: Path,
    *,
    source_backed: bool,
    name: str | None = None,
) -> str:
    from course_helper.cards import VOCABULARY_VERSION_ID, seed_vocabulary
    from course_helper.domain.common import ActorRef, SourceLocator
    from course_helper.domain.knowledge import (
        CardContentNode,
        ChunkCitation,
        KnowledgeCardVersion,
        TagAssignment,
    )
    from course_helper.domain.sources import ChunkLocator, ExtractedChunk, SourceAssetVersion

    now = datetime(2026, 7, 16, tzinfo=timezone.utc)
    actor = ActorRef(actor_type="service", actor_id="api-publish-tests")
    suffix = name or ("valid" if source_backed else "invalid")
    card_version_id = f"publish-candidate-{suffix}"
    with KnowledgeCatalog.open(database) as catalog:
        seed_vocabulary(catalog)
        citations: tuple[ChunkCitation, ...] = ()
        if source_backed:
            source_version_id = f"publish-source-{suffix}"
            chunk_id = f"publish-chunk-{suffix}"
            catalog.insert_source(
                SourceAssetVersion(
                    logical_id=f"logical-{source_version_id}",
                    version_id=source_version_id,
                    revision=1,
                    content_digest=hashlib.sha256(source_version_id.encode()).hexdigest(),
                    created_at=now,
                    created_by=actor,
                    locator=SourceLocator(root_id="fixture", relative_path="publish.md"),
                    display_name="publish.md",
                    source_kind="markdown",
                    media_type="text/markdown",
                    byte_size=10,
                    extraction_status="parsed",
                )
            )
            catalog.insert_chunk(
                ExtractedChunk(
                    chunk_id=chunk_id,
                    source_version_id=source_version_id,
                    ordinal=0,
                    modality="text",
                    language="en",
                    normalized_text="Verified publication evidence",
                    content_digest=hashlib.sha256(b"Verified publication evidence").hexdigest(),
                    locator=ChunkLocator(kind="markdown-section", ast_path=(1,)),
                )
            )
            citations = (
                ChunkCitation(
                    chunk_id=chunk_id,
                    source_version_id=source_version_id,
                    quoted_text="Verified publication evidence",
                ),
            )
        catalog.insert_card(
            KnowledgeCardVersion(
                logical_id=f"logical-{card_version_id}",
                version_id=card_version_id,
                revision=1,
                content_digest=hashlib.sha256(card_version_id.encode()).hexdigest(),
                created_at=now,
                created_by=actor,
                main_type_id="concept",
                title="RFM publication",
                learning_objective="Explain a governed RFM workflow",
                content_ast=(
                    CardContentNode(type="paragraph", text="Governed content"),
                ),
                suggested_minutes=5,
                vocabulary_version_id=VOCABULARY_VERSION_ID,
                tag_assignments=(
                    TagAssignment(
                        vocabulary_version_id=VOCABULARY_VERSION_ID,
                        dimension_id="difficulty",
                        tag_id="difficulty:beginner",
                    ),
                ),
                chunk_citations=citations,
                status="review",
            )
        )
    return card_version_id


def test_real_spawn_publishes_only_an_existing_card_version_id(tmp_path: Path) -> None:
    from course_helper.jobs import BoundedJobRunner, KnowledgePublishJob, WorkerRuntimeConfig

    database = tmp_path / "knowledge.db"
    card_version_id = _seed_publish_candidate(database, source_backed=True)
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        )
    )
    job = KnowledgePublishJob.model_validate(
        {"type": "knowledge_publish", "cardVersionId": card_version_id}
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 200
    assert outcome.evidence.kind == "publish"
    assert outcome.result["submittedCardVersionId"] == card_version_id
    with KnowledgeCatalog.open(database) as catalog:
        status = catalog.connection.execute(
            "SELECT status FROM cards WHERE logical_id = ? ORDER BY revision DESC LIMIT 1",
            (f"logical-{card_version_id}",),
        ).fetchone()[0]
    assert status == "published"


def test_concurrent_publish_jobs_return_only_their_exact_evidence(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from course_helper import jobs as jobs_module
    from course_helper.jobs import KnowledgePublishJob, WorkerRuntimeConfig, _run_knowledge_publish

    database = tmp_path / "knowledge.db"
    card_ids = (
        _seed_publish_candidate(database, source_backed=True, name="concurrent-a"),
        _seed_publish_candidate(database, source_backed=True, name="concurrent-b"),
    )
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(),
    )
    real_publish_card = jobs_module.publish_card
    expected_evidence_ids: dict[str, str] | None = None

    for _round in range(20):
        published = threading.Barrier(2)

        def synchronized_publish(card: object, catalog: KnowledgeCatalog) -> object:
            result = real_publish_card(card, catalog)  # type: ignore[arg-type]
            published.wait(timeout=10)
            return result

        monkeypatch.setattr(jobs_module, "publish_card", synchronized_publish)

        def worker(card_version_id: str) -> tuple[str, dict[str, object], EvidenceObject]:
            result, evidence = _run_knowledge_publish(
                KnowledgePublishJob.model_validate(
                    {
                        "type": "knowledge_publish",
                        "cardVersionId": card_version_id,
                    }
                ),
                config,
            )
            return card_version_id, result, evidence

        with ThreadPoolExecutor(max_workers=2) as executor:
            outcomes = tuple(executor.map(worker, card_ids))

        current_evidence_ids: dict[str, str] = {}
        for submitted_id, result, evidence in outcomes:
            assert result["submittedCardVersionId"] == submitted_id
            assert evidence.kind == "publish"
            assert evidence.subject_version_id == result["publishedCardVersionId"]
            current_evidence_ids[submitted_id] = evidence.evidence_id
        assert len(set(current_evidence_ids.values())) == 2
        if expected_evidence_ids is None:
            expected_evidence_ids = current_evidence_ids
        else:
            assert current_evidence_ids == expected_evidence_ids


def test_dedup_publish_retry_returns_the_original_archived_subject_evidence(
    tmp_path: Path,
) -> None:
    from course_helper.jobs import KnowledgePublishJob, WorkerRuntimeConfig, _run_knowledge_publish

    database = tmp_path / "knowledge.db"
    original_id = _seed_publish_candidate(database, source_backed=True, name="dedup-original")
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(),
    )
    original_result, _original_evidence = _run_knowledge_publish(
        KnowledgePublishJob.model_validate(
            {"type": "knowledge_publish", "cardVersionId": original_id}
        ),
        config,
    )
    duplicate_id = "publish-candidate-dedup-copy"
    with KnowledgeCatalog.open(database) as catalog:
        original_row = catalog.connection.execute(
            "SELECT payload_json FROM cards WHERE version_id = ?",
            (original_id,),
        ).fetchone()
        assert original_row is not None
        from course_helper.domain.knowledge import KnowledgeCardVersion

        duplicate = KnowledgeCardVersion.model_validate_json(original_row[0]).model_copy(
            update={
                "logical_id": "logical-publish-candidate-dedup-copy",
                "version_id": duplicate_id,
                "content_digest": hashlib.sha256(duplicate_id.encode()).hexdigest(),
            }
        )
        catalog.insert_card(duplicate)

    first_result, first_evidence = _run_knowledge_publish(
        KnowledgePublishJob.model_validate(
            {"type": "knowledge_publish", "cardVersionId": duplicate_id}
        ),
        config,
    )
    unrelated_id = _seed_publish_candidate(database, source_backed=True, name="after-dedup")
    _run_knowledge_publish(
        KnowledgePublishJob.model_validate(
            {"type": "knowledge_publish", "cardVersionId": unrelated_id}
        ),
        config,
    )
    retry_result, retry_evidence = _run_knowledge_publish(
        KnowledgePublishJob.model_validate(
            {"type": "knowledge_publish", "cardVersionId": duplicate_id}
        ),
        config,
    )

    assert first_result == retry_result
    assert first_result["publishedCardVersionId"] == original_result["publishedCardVersionId"]
    assert first_evidence.kind == retry_evidence.kind == "dedup"
    assert retry_evidence.evidence_id == first_evidence.evidence_id
    assert retry_evidence.subject_version_id == first_evidence.subject_version_id
    assert retry_evidence.subject_version_id not in {duplicate_id, first_result["publishedCardVersionId"]}
    with KnowledgeCatalog.open(database) as catalog:
        lineage = catalog.connection.execute(
            "SELECT to_version_id, evidence_id FROM lineage "
            "WHERE from_version_id = ? AND relation = 'deduplicates'",
            (retry_evidence.subject_version_id,),
        ).fetchone()
    assert lineage == (first_result["publishedCardVersionId"], first_evidence.evidence_id)


def test_real_spawn_publish_failure_is_sanitized_and_keeps_review_state(
    tmp_path: Path,
) -> None:
    from course_helper.jobs import BoundedJobRunner, KnowledgePublishJob, WorkerRuntimeConfig

    database = tmp_path / "knowledge.db"
    card_version_id = _seed_publish_candidate(database, source_backed=False)
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        )
    )
    job = KnowledgePublishJob.model_validate(
        {"type": "knowledge_publish", "cardVersionId": card_version_id}
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 422
    assert outcome.evidence.checks[0].code == "publish-blocked"
    assert str(tmp_path) not in outcome.evidence.model_dump_json()
    with KnowledgeCatalog.open(database) as catalog:
        stored = catalog.connection.execute(
            "SELECT status FROM cards WHERE version_id = ?",
            (card_version_id,),
        ).fetchone()
        evidence_count = catalog.connection.execute(
            "SELECT count(*) FROM evidence"
        ).fetchone()[0]
    assert stored == ("review",)
    assert evidence_count == 0


def test_missing_publish_card_is_422_without_spawning(tmp_path: Path) -> None:
    from course_helper.jobs import BoundedJobRunner, KnowledgePublishJob, WorkerRuntimeConfig

    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    context = _FakeContext()
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(),
        ),
        mp_context=context,
    )
    job = KnowledgePublishJob.model_validate(
        {"type": "knowledge_publish", "cardVersionId": "missing-card-version"}
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 422
    assert context.process_calls == 0
    assert outcome.evidence.checks[0].code == "job-preflight"
    assert outcome.evidence.checks[0].details["reason_code"] == "card-not-found"


def test_real_spawn_maps_immutable_catalog_conflict_to_stable_422(
    tmp_path: Path,
) -> None:
    from course_helper.domain.common import SourceLocator
    from course_helper.jobs import BoundedJobRunner, SourceIngestJob, WorkerRuntimeConfig
    from course_helper.parsers.markdown_parser import MarkdownParser
    from course_helper.source_roots import SourceRootRegistry

    source_root = tmp_path / "sources"
    source_root.mkdir()
    (source_root / "lesson.md").write_text("# Lesson\n\nStable body.\n", encoding="utf-8")
    registry = SourceRootRegistry({"fixture": source_root})
    extraction = MarkdownParser(registry).parse(
        SourceLocator(root_id="fixture", relative_path="lesson.md")
    )
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database) as catalog:
        catalog.insert_source(
            extraction.source.model_copy(update={"display_name": "conflicting-name.md"})
        )
    runner = BoundedJobRunner(
        WorkerRuntimeConfig(
            database_path=str(database),
            app_data_path=str(tmp_path / "app-data"),
            source_roots=(("fixture", str(source_root)),),
        )
    )
    job = SourceIngestJob.model_validate(
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "lesson.md"},
            "selection": {"headingSelectors": []},
        }
    )

    outcome = asyncio.run(runner.run(job, disconnected=_connected))

    assert outcome.status_code == 422
    assert outcome.evidence.checks[0].code == "immutable-version-conflict"
    assert str(tmp_path) not in outcome.evidence.model_dump_json()
    with KnowledgeCatalog.open(database) as catalog:
        assert catalog.connection.execute("SELECT count(*) FROM sources").fetchone()[0] == 1
        assert catalog.connection.execute("SELECT count(*) FROM chunks").fetchone()[0] == 0
        assert catalog.connection.execute("SELECT count(*) FROM evidence").fetchone()[0] == 0


def test_child_rechecks_dataset_size_after_parent_preflight(tmp_path: Path) -> None:
    from course_helper.jobs import (
        BoundedJobRunner,
        DatasetProfileJob,
        WorkerRuntimeConfig,
        _spawn_job_entry,
    )

    source_root = tmp_path / "sources"
    source_root.mkdir()
    dataset = source_root / "mutable.parquet"
    dataset.write_bytes(b"PAR1")
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(("fixture", str(source_root)),),
    )
    runner = BoundedJobRunner(config, mp_context=_FakeContext())
    job = DatasetProfileJob.model_validate(
        {
            "type": "dataset_profile",
            "locator": {"rootId": "fixture", "relativePath": "mutable.parquet"},
            "sampleLimit": 20,
        }
    )
    assert runner._preflight(job)["timeoutSeconds"] == 60
    with dataset.open("r+b") as output:
        output.truncate(1024 * 1024 * 1024 + 1)
    result_queue = _RecordingQueue()

    _spawn_job_entry(
        job.model_dump(mode="json", by_alias=True),
        config,
        result_queue,
    )

    assert result_queue.items == [
        {
            "ok": False,
            "statusCode": 413,
            "code": "job-preflight",
            "reasonCode": "dataset-too-large",
        }
    ]
    with KnowledgeCatalog.open(database) as catalog:
        assert catalog.connection.execute("SELECT count(*) FROM datasets").fetchone()[0] == 0


def test_child_rechecks_pptx_slide_count_after_parent_preflight(tmp_path: Path) -> None:
    from pptx import Presentation

    from course_helper.jobs import (
        BoundedJobRunner,
        SourceIngestJob,
        WorkerRuntimeConfig,
        _spawn_job_entry,
    )

    source_root = tmp_path / "sources"
    source_root.mkdir()
    presentation_path = source_root / "mutable.pptx"

    def write_slides(count: int) -> None:
        presentation = Presentation()
        for _ in range(count):
            presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(presentation_path)

    write_slides(1)
    database = tmp_path / "knowledge.db"
    with KnowledgeCatalog.open(database):
        pass
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(("fixture", str(source_root)),),
    )
    runner = BoundedJobRunner(config, mp_context=_FakeContext())
    job = SourceIngestJob.model_validate(
        {
            "type": "source_ingest",
            "locator": {"rootId": "fixture", "relativePath": "mutable.pptx"},
            "selection": {"slideNumbers": []},
        }
    )
    assert runner._preflight(job)["maxSelectedSlides"] == 64
    write_slides(65)
    result_queue = _RecordingQueue()

    _spawn_job_entry(
        job.model_dump(mode="json", by_alias=True),
        config,
        result_queue,
    )

    assert result_queue.items == [
        {
            "ok": False,
            "statusCode": 422,
            "code": "job-preflight",
            "reasonCode": "too-many-slides",
        }
    ]
    with KnowledgeCatalog.open(database) as catalog:
        assert catalog.connection.execute("SELECT count(*) FROM sources").fetchone()[0] == 0
        assert catalog.connection.execute("SELECT count(*) FROM evidence").fetchone()[0] == 0


def test_child_rechecks_card_existence_after_parent_preflight(tmp_path: Path) -> None:
    from course_helper.jobs import (
        BoundedJobRunner,
        KnowledgePublishJob,
        WorkerRuntimeConfig,
        _spawn_job_entry,
    )

    database = tmp_path / "knowledge.db"
    card_version_id = _seed_publish_candidate(database, source_backed=True)
    config = WorkerRuntimeConfig(
        database_path=str(database),
        app_data_path=str(tmp_path / "app-data"),
        source_roots=(),
    )
    runner = BoundedJobRunner(config, mp_context=_FakeContext())
    job = KnowledgePublishJob.model_validate(
        {"type": "knowledge_publish", "cardVersionId": card_version_id}
    )
    assert runner._preflight(job)["maxCardVersions"] == 1
    with KnowledgeCatalog.open(database) as catalog:
        # Test-only fault injection: preserve production's restrictive audit FKs and
        # explicitly remove lifecycle state before simulating a missing raw card.
        catalog.connection.execute(
            "DROP TRIGGER card_lifecycle_events_append_only_delete"
        )
        catalog.connection.execute(
            "DELETE FROM card_lifecycle_current WHERE card_version_id = ?",
            (card_version_id,),
        )
        catalog.connection.execute(
            "DELETE FROM card_lifecycle_events WHERE card_version_id = ?",
            (card_version_id,),
        )
        catalog.connection.execute(
            "DELETE FROM card_tags WHERE card_version_id = ?",
            (card_version_id,),
        )
        catalog.connection.execute(
            "DELETE FROM cards WHERE version_id = ?",
            (card_version_id,),
        )
        catalog.connection.commit()
    result_queue = _RecordingQueue()

    _spawn_job_entry(
        job.model_dump(mode="json", by_alias=True),
        config,
        result_queue,
    )

    assert result_queue.items == [
        {
            "ok": False,
            "statusCode": 422,
            "code": "job-preflight",
            "reasonCode": "card-not-found",
        }
    ]
    with KnowledgeCatalog.open(database) as catalog:
        assert catalog.connection.execute("SELECT count(*) FROM evidence").fetchone()[0] == 0
